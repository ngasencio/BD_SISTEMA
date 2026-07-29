"""
Módulo PAC — generación de reportería institucional (Word, PPT, PDF) para el
Servicio de Salud Osorno. Informe único y exhaustivo (decisión del usuario
2026-07-21): consolida TODA la información del dashboard /pac-cumplimiento
(Resumen, Jerarquía, Rankings, Cumplimiento Temporal FSC, y Ejecución del Plan
de Compras vía Ficha PAC/OC) en un solo documento por formato — capítulos por
subdirección, portada institucional (logo + fachada), índice con numeración de
página real, y pie de página en todas las hojas.

PPT es deliberadamente más breve (versión ejecutiva pensada para presentar en
una reunión de Alta Dirección) — Word/PDF llevan el detalle exhaustivo con
tablas por departamento.

Paleta institucional: mismos colores que frontend/src/index.css
(--gob-azul/--gob-celeste) y que el dashboard (verde=Dentro/Ejecutado,
rojo=Fuera/Atrasado, ámbar=Pendiente, gris=Sin dato).
"""
import io
import os
from collections import defaultdict
from datetime import date

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from PIL import Image as PILImage

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch, cm
from reportlab.platypus import (
    BaseDocTemplate, PageTemplate, Frame, Paragraph, Spacer, Table, TableStyle,
    Image as RLImage, PageBreak,
)
from reportlab.platypus.tableofcontents import TableOfContents

from .plantillas_narrativas import (
    etiqueta_periodo, parrafo_resumen_ejecutivo, parrafo_cumplimiento_temporal,
    parrafo_capitulo_subdireccion, parrafo_conclusiones, _n, _money,
)

# =============================================================================
# Identidad institucional
# =============================================================================

NOMBRE_INSTITUCION = 'Dirección Servicio de Salud Osorno'
TITULO_INFORME = 'Informe de Seguimiento y Ejecución del Plan Anual de Compras'

COLOR_INSTITUCIONAL = '#1e3a5f'        # --gob-azul (frontend/src/index.css)
COLOR_INSTITUCIONAL_CLARO = '#38b2bd'  # --gob-celeste
COLOR_DENTRO = '#16a34a'
COLOR_FUERA = '#dc2626'
COLOR_PENDIENTE = '#d97706'
COLOR_SIN_DATO = '#9ca3af'

NOMBRES_SUBDIRECCIONES_INSTITUCIONALES = {
    'DIRECTOR',
    'SUBDIRECCION DE GESTION ASISTENCIAL',
    'SUBDIRECCION ADMINISTRATIVA',
    'SUBDIRECCION DE GESTION Y DESARROLLO DE LAS PERSONAS',
}

# Los nombres de subdirección llegan en mayúsculas y sin tildes desde el origen
# (`SsoSubdireccion.nombre`/`PlanerPAC.sub`) — un informe institucional para Alta
# Dirección merece la ortografía correcta al menos para estas 4 ramas fijas, que
# son además los títulos de capítulo más visibles del documento (encabezados +
# índice). No se intenta corregir los ~150 nombres de departamento (serían
# demasiados para mantener a mano) — solo las subdirecciones.
_NOMBRES_DISPLAY_SUBDIRECCION = {
    'DIRECTOR': 'Director',
    'SUBDIRECCION DE GESTION ASISTENCIAL': 'Subdirección de Gestión Asistencial',
    'SUBDIRECCION ADMINISTRATIVA': 'Subdirección Administrativa',
    'SUBDIRECCION DE GESTION Y DESARROLLO DE LAS PERSONAS': 'Subdirección de Gestión y Desarrollo de las Personas',
    'SIN CLASIFICAR': 'Sin Clasificar',
}


def _nombre_subdireccion_display(nombre):
    if not nombre:
        return nombre
    return _NOMBRES_DISPLAY_SUBDIRECCION.get(nombre.upper(), nombre.title())


_ORDEN_SUBDIRECCIONES_INSTITUCIONALES = [
    'DIRECTOR', 'SUBDIRECCION DE GESTION ASISTENCIAL',
    'SUBDIRECCION ADMINISTRATIVA', 'SUBDIRECCION DE GESTION Y DESARROLLO DE LAS PERSONAS',
]


def _agrupar_por_subdireccion(items, campo='subdireccion_nombre'):
    """Agrupa una lista plana de fichas/proyectos por nombre de subdirección — usado
    por 'Alertas y Seguimiento' para que 'Fichas a ejecutar el próximo mes', 'Fichas
    atrasadas' y 'Proyectos sin iniciar' se lean por subdirección en vez de una sola
    tabla mezclada (pedido del usuario 2026-07-23, para mayor claridad). Orden: las 4
    subdirecciones institucionales primero (orden fijo), luego otras ramas de la red
    en alfabético, 'Sin Clasificar' siempre al final."""
    por_nombre = defaultdict(list)
    for item in items:
        por_nombre[item.get(campo) or 'Sin Clasificar'].append(item)

    def _orden(nombre):
        if nombre == 'Sin Clasificar':
            return (2, nombre)
        if nombre in _ORDEN_SUBDIRECCIONES_INSTITUCIONALES:
            return (0, _ORDEN_SUBDIRECCIONES_INSTITUCIONALES.index(nombre))
        return (1, nombre)

    return [
        {'nombre': nombre, 'nombre_display': _nombre_subdireccion_display(nombre), 'items': por_nombre[nombre]}
        for nombre in sorted(por_nombre, key=_orden)
    ]


def _total_items_grupos(grupos):
    return sum(len(g['items']) for g in grupos)


_BASE_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
_LOGO_PATH = os.path.join(_BASE_DIR, 'frontend', 'public', 'logo.jpg')
_EDIFICIO_PATH = os.path.join(_BASE_DIR, 'frontend', 'public', 'edificio.jpg')

_CACHE_IMAGENES = {}


def _imagen_redimensionada(path, ancho_max):
    """Carga y cachea (en memoria del proceso) una versión redimensionada de una
    imagen institucional — `logo.jpg` pesa 3.6MB a 5906x5349px; insertarlo tal cual
    en cada reporte generado infla el archivo de salida y ralentiza la generación."""
    clave = (path, ancho_max)
    if clave not in _CACHE_IMAGENES:
        try:
            img = PILImage.open(path).convert('RGB')
            if img.width > ancho_max:
                ratio = ancho_max / img.width
                img = img.resize((ancho_max, int(img.height * ratio)), PILImage.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format='JPEG', quality=87)
            _CACHE_IMAGENES[clave] = buf.getvalue()
        except Exception:
            _CACHE_IMAGENES[clave] = None
    data = _CACHE_IMAGENES[clave]
    return io.BytesIO(data) if data else None


def _logo_bytes(ancho_max=500):
    return _imagen_redimensionada(_LOGO_PATH, ancho_max)


def _edificio_bytes(ancho_max=1000):
    return _imagen_redimensionada(_EDIFICIO_PATH, ancho_max)


plt.rcParams.update({
    'font.family': 'sans-serif',
    'font.size': 10,
    'axes.edgecolor': '#cbd5e1',
    'axes.labelcolor': '#374151',
    'text.color': '#374151',
    'xtick.color': '#64748b',
    'ytick.color': '#64748b',
})


def _fig_a_bytes(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight', transparent=False, facecolor='white')
    plt.close(fig)
    buf.seek(0)
    return buf


def _truncar(texto, largo):
    texto = str(texto or '—')
    return texto if len(texto) <= largo else texto[:largo - 1] + '…'


# =============================================================================
# Gráficos (matplotlib → PNG) — reutilizados por Word/PPT/PDF
# =============================================================================

def grafico_donut_dentro_fuera(pct_dentro, titulo='Dentro / Fuera PAC'):
    """Dona simple 2 segmentos — % Dentro/Fuera PAC a nivel de formularios."""
    pct_fuera = round(100 - pct_dentro, 1)
    fig, ax = plt.subplots(figsize=(3.6, 3.6))
    ax.pie(
        [pct_dentro, pct_fuera],
        colors=[COLOR_DENTRO, COLOR_FUERA],
        startangle=90,
        wedgeprops={'width': 0.35, 'edgecolor': 'white', 'linewidth': 2},
        autopct=lambda p: f'{p:.0f}%' if p > 5 else '',
        pctdistance=0.82,
        textprops={'color': 'white', 'fontweight': 'bold', 'fontsize': 11},
    )
    ax.set_title(titulo, fontsize=11, pad=10)
    fig.legend(
        ['Dentro PAC', 'Fuera PAC'], loc='lower center', ncol=2, frameon=False,
        bbox_to_anchor=(0.5, -0.02), fontsize=9,
    )
    return _fig_a_bytes(fig)


def grafico_donut_ejecucion_planer(ejecutados, pendientes, atrasados, titulo='Ejecución Plan de Compras'):
    """Dona 3 segmentos — Ejecutado/Pendiente/Atrasado, dominio Ficha PAC (`PlanerPAC`).
    Distinto de `grafico_donut_cumplimiento_temporal` (dominio formularios FSC) — no
    reutilizarlo acá evita mostrar la etiqueta 'En fecha' donde correspondía 'Ejecutado'."""
    valores = [ejecutados, pendientes, atrasados]
    if not sum(valores):
        return None
    etiquetas = ['Ejecutado', 'Pendiente', 'Atrasado']
    colores = [COLOR_DENTRO, COLOR_PENDIENTE, COLOR_FUERA]
    datos = [(e, v, c) for e, v, c in zip(etiquetas, valores, colores) if v > 0]

    fig, ax = plt.subplots(figsize=(4.2, 3.8))
    ax.pie(
        [d[1] for d in datos], colors=[d[2] for d in datos], startangle=90,
        wedgeprops={'width': 0.35, 'edgecolor': 'white', 'linewidth': 2},
        autopct=lambda p: f'{p:.0f}%' if p > 5 else '',
        pctdistance=0.82,
        textprops={'color': 'white', 'fontweight': 'bold', 'fontsize': 10},
    )
    ax.set_title(titulo, fontsize=11, pad=10)
    fig.legend(
        [d[0] for d in datos], loc='lower center', ncol=3, frameon=False,
        bbox_to_anchor=(0.5, -0.02), fontsize=8,
    )
    return _fig_a_bytes(fig)


def grafico_barras_comparativa_anual(comparativa_anual, titulo='Comparativa histórica — Dentro vs Fuera PAC'):
    """Barras apiladas por año — comparativa_anual: [{anho, dentro, fuera, ...}]."""
    if not comparativa_anual:
        return None
    anhos = [str(r['anho']) for r in comparativa_anual]
    dentro = [r['dentro'] for r in comparativa_anual]
    fuera = [r['fuera'] for r in comparativa_anual]

    fig, ax = plt.subplots(figsize=(6.4, 3.6))
    ax.bar(anhos, dentro, color=COLOR_DENTRO, label='Dentro PAC', width=0.55)
    ax.bar(anhos, fuera, bottom=dentro, color=COLOR_FUERA, label='Fuera PAC', width=0.55)
    ax.set_title(titulo, fontsize=11, pad=10)
    ax.spines[['top', 'right']].set_visible(False)
    ax.legend(loc='upper center', ncol=2, frameon=False, bbox_to_anchor=(0.5, -0.12))
    fig.tight_layout()
    return _fig_a_bytes(fig)


def grafico_donut_cumplimiento_temporal(kpis_temporal, titulo='Cumplimiento Temporal'):
    """Dona 4 segmentos — En fecha / Atrasado / Pendiente / Sin planificación con fecha."""
    valores = [
        kpis_temporal.get('en_fecha', 0),
        kpis_temporal.get('atrasado', 0),
        kpis_temporal.get('pendiente', 0),
        kpis_temporal.get('sin_planificacion_con_fecha', 0),
    ]
    if not sum(valores):
        return None
    etiquetas = ['En fecha', 'Atrasado', 'Pendiente', 'Sin planificación']
    colores = [COLOR_DENTRO, COLOR_FUERA, COLOR_PENDIENTE, COLOR_SIN_DATO]
    datos = [(e, v, c) for e, v, c in zip(etiquetas, valores, colores) if v > 0]

    fig, ax = plt.subplots(figsize=(4.2, 3.8))
    ax.pie(
        [d[1] for d in datos], colors=[d[2] for d in datos], startangle=90,
        wedgeprops={'width': 0.35, 'edgecolor': 'white', 'linewidth': 2},
    )
    ax.set_title(titulo, fontsize=11, pad=10)
    fig.legend(
        [d[0] for d in datos], loc='lower center', ncol=2, frameon=False,
        bbox_to_anchor=(0.5, -0.05), fontsize=8,
    )
    return _fig_a_bytes(fig)


def grafico_barras_ranking(filas, campo_nombre='nombre', campo_valor='score', titulo='Ranking', color=COLOR_DENTRO, horizontal=True):
    """Barras horizontales para rankings mejores/peores (departamento)."""
    if not filas:
        return None
    nombres = [_truncar(f[campo_nombre], 40) for f in filas]
    valores = [f[campo_valor] for f in filas]

    fig, ax = plt.subplots(figsize=(6.4, max(2.2, 0.35 * len(filas))))
    if horizontal:
        ax.barh(nombres[::-1], valores[::-1], color=color, height=0.6)
        ax.set_xlabel('Score')
    else:
        ax.bar(nombres, valores, color=color, width=0.6)
        ax.set_ylabel('Score')
        plt.xticks(rotation=45, ha='right')
    ax.set_title(titulo, fontsize=11, pad=10)
    ax.spines[['top', 'right']].set_visible(False)
    fig.tight_layout()
    return _fig_a_bytes(fig)


def grafico_evolucion_mensual_dentro_fuera(meses, anho, anho_anterior, titulo=None):
    """Barras Dentro/Fuera PAC por mes de `anho` + línea de % Dentro del año
    anterior superpuesta — mismo dato que el panel 'Temporalidad' del Resumen."""
    if not meses or not any(m['total'] for m in meses):
        return None
    nombres_mes = [m['nombre_mes'] for m in meses]
    dentro = [m['dentro'] for m in meses]
    fuera = [m['fuera'] for m in meses]
    pct_anterior = [m['pct_dentro_anho_anterior'] for m in meses]

    fig, ax1 = plt.subplots(figsize=(7.2, 3.3))
    ax1.bar(nombres_mes, dentro, color=COLOR_DENTRO, label=f'Dentro PAC {anho}')
    ax1.bar(nombres_mes, fuera, bottom=dentro, color=COLOR_FUERA, label=f'Fuera PAC {anho}')
    ax1.set_ylabel('N° formularios')
    ax1.spines[['top']].set_visible(False)

    ax2 = ax1.twinx()
    puntos = [(i, v) for i, v in enumerate(pct_anterior) if v is not None]
    if puntos:
        xs, ys = zip(*puntos)
        ax2.plot(xs, ys, color='#64748b', linestyle='--', marker='o', markersize=3, label=f'% Dentro {anho_anterior}')
    ax2.set_ylabel(f'% Dentro {anho_anterior}')
    ax2.set_ylim(0, 100)

    ax1.set_title(titulo or 'Evolución mensual — Dentro/Fuera PAC', fontsize=11, pad=10)
    l1, e1 = ax1.get_legend_handles_labels()
    l2, e2 = ax2.get_legend_handles_labels()
    ax1.legend(l1 + l2, e1 + e2, loc='upper center', bbox_to_anchor=(0.5, -0.16), ncol=3, fontsize=8, frameon=False)
    fig.tight_layout()
    return _fig_a_bytes(fig)


def grafico_barras_ejecucion_mensual_planer(meses, titulo=None):
    """Barras apiladas Ejecutado/Pendiente/Atrasado por mes — Ejecución del Plan de
    Compras (Ficha PAC), mismo dato que el gráfico del sub-tab Temporal."""
    if not meses or not any(m['total'] for m in meses):
        return None
    nombres_mes = [m['nombre_mes'] for m in meses]
    ejecutados = [m['ejecutados'] for m in meses]
    pendientes = [m['pendientes'] for m in meses]
    atrasados = [m['atrasados'] for m in meses]
    base_atrasado = [e + p for e, p in zip(ejecutados, pendientes)]

    fig, ax = plt.subplots(figsize=(7.2, 3))
    ax.bar(nombres_mes, ejecutados, color=COLOR_DENTRO, label='Ejecutado')
    ax.bar(nombres_mes, pendientes, bottom=ejecutados, color=COLOR_PENDIENTE, label='Pendiente')
    ax.bar(nombres_mes, atrasados, bottom=base_atrasado, color=COLOR_FUERA, label='Atrasado')
    ax.set_ylabel('N° fichas')
    ax.spines[['top']].set_visible(False)
    ax.set_title(titulo or 'Ejecución mensual del Plan de Compras', fontsize=11, pad=10)
    ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.14), ncol=3, fontsize=8, frameon=False)
    fig.tight_layout()
    return _fig_a_bytes(fig)


def grafico_barras_ejecucion_por_depto(departamentos, titulo=None):
    """Barras horizontales apiladas Ejecutado/Pendiente/Atrasado por departamento —
    versión del gráfico de `grafico_barras_ejecucion_mensual_planer` pero con eje
    departamento en vez de mes, para el capítulo por subdirección del PPT."""
    filas = [d for d in departamentos if d['total']]
    if not filas:
        return None
    nombres = [_truncar(d['nombre'].title(), 32) for d in filas]
    ejecutados = [d['ejecutados'] for d in filas]
    pendientes = [d['pendientes'] for d in filas]
    atrasados = [d['atrasados'] for d in filas]
    base_atrasado = [e + p for e, p in zip(ejecutados, pendientes)]

    fig, ax = plt.subplots(figsize=(6.4, max(2.2, 0.45 * len(filas))))
    ax.barh(nombres[::-1], ejecutados[::-1], color=COLOR_DENTRO, label='Ejecutado')
    ax.barh(nombres[::-1], pendientes[::-1], left=ejecutados[::-1], color=COLOR_PENDIENTE, label='Pendiente')
    ax.barh(nombres[::-1], atrasados[::-1], left=base_atrasado[::-1], color=COLOR_FUERA, label='Atrasado')
    ax.set_xlabel('N° fichas')
    ax.spines[['top', 'right']].set_visible(False)
    ax.set_title(titulo or 'Ejecución del Plan de Compras por departamento', fontsize=11, pad=10)
    ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.12), ncol=3, fontsize=8, frameon=False)
    fig.tight_layout()
    return _fig_a_bytes(fig)


# =============================================================================
# Recolección y combinación de datos — TODO el dashboard /pac-cumplimiento en
# una sola función, para que los 3 formatos consuman exactamente lo mismo.
# =============================================================================

def _anho_de_periodo(periodo):
    return int(periodo[:4])


def _datos_informe_completo(periodo):
    """Reúne TODA la información del dashboard PAC Cumplimiento para un `periodo`
    ('YYYY-MM' o 'YYYY-QN'): las secciones basadas en FormularioFSCDerivado respetan
    el período exacto (mes/trimestre); la sección de Ejecución del Plan de Compras
    (Ficha PAC/OC) siempre refleja el AÑO PAC completo extraído de `periodo`, porque
    esa planificación es intrínsecamente anual — se etiqueta así explícitamente en
    el informe para no confundir ambos cortes temporales."""
    from .services import (
        calcular_pac_dentro_fuera_stats, calcular_pac_comparativa_periodos,
        calcular_pac_cumplimiento_temporal, calcular_pac_jerarquia, calcular_pac_rankings,
        calcular_pac_resumen_subdireccion, calcular_pac_temporalidad_mensual,
        calcular_pac_jerarquia_planer, calcular_pac_temporal_mensual_planer,
        calcular_pac_detalle_fsc_por_subdireccion,
        _calcular_fichas_pac_completo, _rango_fechas_periodo,
    )

    anho = _anho_de_periodo(periodo)
    desde, hasta = _rango_fechas_periodo(periodo)
    desde_iso, hasta_iso = desde.isoformat(), hasta.isoformat()
    label = etiqueta_periodo(periodo)

    dentro_fuera = calcular_pac_dentro_fuera_stats(fecha_desde=desde_iso, fecha_hasta=hasta_iso)
    comparativa_periodos = calcular_pac_comparativa_periodos(periodo)
    # `anho=anho` es obligatorio acá aunque fecha_desde/fecha_hasta ya acoten el período: las
    # 3 funciones de abajo usan `anho` para escoger los eventos planificados de PlanerPAC
    # contra los que se compara "% en fecha" (`_eventos_planificados_por_proyecto`). Sin
    # `anho`, ese universo de eventos queda SIN filtrar por año (todo el histórico PC20-PC26+
    # en vez de solo el año del período), dando un "% en fecha"/score de ranking distinto al
    # que muestra el dashboard interactivo para el mismo filtro — bug real encontrado en
    # revisión de código 2026-07-22, coherente con cómo la vista del dashboard SÍ pasa `anho`
    # (ver `pac_cumplimiento_temporal_view` en views.py).
    temporal = calcular_pac_cumplimiento_temporal(anho=anho, fecha_desde=desde_iso, fecha_hasta=hasta_iso)
    jerarquia_fsc = calcular_pac_jerarquia(anho=anho, fecha_desde=desde_iso, fecha_hasta=hasta_iso)
    rankings_depto = calcular_pac_rankings(anho=anho, fecha_desde=desde_iso, fecha_hasta=hasta_iso, tipo='depto')
    rankings_formulario = calcular_pac_rankings(anho=anho, fecha_desde=desde_iso, fecha_hasta=hasta_iso, tipo='formulario')
    resumen_subdireccion = calcular_pac_resumen_subdireccion(anho)
    temporalidad_mensual = calcular_pac_temporalidad_mensual(anho)

    fichas = _calcular_fichas_pac_completo(anho=anho)
    jerarquia_planer = calcular_pac_jerarquia_planer(anho)
    temporal_mensual_planer = calcular_pac_temporal_mensual_planer(anho)
    # Detalle FSC individual por subdirección (folio, unidad, monto, fecha) — complementa
    # `jerarquia_fsc` (que solo agrega por departamento) para poder listar, dentro del
    # capítulo de cada subdirección, CON QUÉ formularios concretos se compone su cifra.
    detalle_fsc_subdireccion = calcular_pac_detalle_fsc_por_subdireccion(
        anho=anho, fecha_desde=desde_iso, fecha_hasta=hasta_iso,
    )

    hoy = date.today()
    anho_prox, mes_prox = (hoy.year, hoy.month + 1) if hoy.month < 12 else (hoy.year + 1, 1)
    prefijo_prox = f'{anho_prox}-{mes_prox:02d}'
    proximo_mes = [
        f for f in fichas
        if f['estado_ejecucion'] != 'EJECUTADO' and (f['fecha_mas_proxima'] or '').startswith(prefijo_prox)
    ]
    atrasadas = [f for f in fichas if f['estado_ejecucion'] == 'ATRASADO']
    total_fichas = len(fichas)
    ejecutados_fichas = sum(1 for f in fichas if f['estado_ejecucion'] == 'EJECUTADO')
    monto_total_fichas = sum(f['monto_total'] or 0 for f in fichas)

    # `proyectos_sin_iniciar` (calcular_pac_cumplimiento_temporal) solo trae id_proyecto/fecha —
    # se enriquece acá con nombre y subdirección (ya calculados en `fichas`) para que Alertas y
    # Seguimiento pueda mostrar el nombre del proyecto y agruparlo por subdirección, sin tocar
    # el contrato de esa función (también la consume el dashboard interactivo).
    fichas_por_id_proyecto = {f['id_proyecto']: f for f in fichas}
    for p in temporal['proyectos_sin_iniciar']:
        ficha_ref = fichas_por_id_proyecto.get(p['id_proyecto'])
        p['nombre_proyecto'] = ficha_ref['nombre_proyecto'] if ficha_ref else None
        p['subdireccion_nombre'] = ficha_ref['subdireccion_nombre'] if ficha_ref else None

    subdirecciones = _combinar_subdirecciones(jerarquia_fsc, jerarquia_planer, fichas, detalle_fsc_subdireccion)
    avance_trimestral = _avance_trimestral(periodo, temporalidad_mensual, temporal_mensual_planer)

    return {
        'periodo': periodo, 'label': label, 'anho': anho, 'anho_anterior': anho - 1, 'hoy': hoy,
        'dentro_fuera': dentro_fuera, 'comparativa_periodos': comparativa_periodos,
        'temporal': temporal, 'jerarquia_fsc': jerarquia_fsc,
        'rankings_depto': rankings_depto, 'rankings_formulario': rankings_formulario,
        'resumen_subdireccion': resumen_subdireccion, 'temporalidad_mensual': temporalidad_mensual,
        'fichas': fichas, 'jerarquia_planer': jerarquia_planer,
        'temporal_mensual_planer': temporal_mensual_planer,
        'proximo_mes': _agrupar_por_subdireccion(proximo_mes),
        'atrasadas': _agrupar_por_subdireccion(atrasadas),
        'proyectos_sin_iniciar': _agrupar_por_subdireccion(temporal['proyectos_sin_iniciar']),
        'total_fichas': total_fichas, 'ejecutados_fichas': ejecutados_fichas,
        'pct_ejecutado_fichas': round(ejecutados_fichas / total_fichas * 100, 1) if total_fichas else 0,
        'monto_total_fichas': monto_total_fichas,
        'subdirecciones': subdirecciones,
        'avance_trimestral': avance_trimestral,
    }


def _avance_trimestral(periodo, temporalidad_mensual, temporal_mensual_planer):
    """Si `periodo` es un trimestre ('YYYY-QN'), arma el avance MES A MES de los 3
    meses que lo componen (Dentro/Fuera PAC + Ejecución del Plan de Compras) — pedido
    explícito del usuario 2026-07-22 ("cuando seleccione trimestral, mostrar avance
    trimestral, para ver su mejora y cómo vamos cumpliendo"). Reutiliza las series de
    12 meses que YA calcula `_datos_informe_completo` (`calcular_pac_temporalidad_mensual`/
    `calcular_pac_temporal_mensual_planer`), solo filtrando a los 3 meses del trimestre
    — no agrega ninguna consulta nueva a la BD. Se aplica a cualquier trimestre (Q1-Q4),
    no solo a Q2/Q3 (los ejemplos que dio el usuario) — no hay razón para mostrar la
    evolución interna solo en algunos trimestres y no en otros. Retorna `None` si
    `periodo` es mensual (un solo mes no tiene "avance interno" que mostrar)."""
    periodo = periodo.upper()
    if 'Q' not in periodo:
        return None
    _, q_str = periodo.split('-Q')
    trimestre = int(q_str)
    mes_ini = (trimestre - 1) * 3 + 1
    meses_trimestre = {mes_ini, mes_ini + 1, mes_ini + 2}

    meses_fsc = [m for m in temporalidad_mensual['meses'] if m['mes'] in meses_trimestre]
    meses_planer = [m for m in temporal_mensual_planer['meses'] if m['mes'] in meses_trimestre]

    pct_dentro_validos = [m['pct_dentro'] for m in meses_fsc if m['pct_dentro'] is not None and m['total']]
    variacion_pct_dentro = round(pct_dentro_validos[-1] - pct_dentro_validos[0], 1) if len(pct_dentro_validos) >= 2 else None

    pct_ejecutado_validos = [m['pct_ejecutado'] for m in meses_planer if m['total']]
    variacion_pct_ejecutado = round(pct_ejecutado_validos[-1] - pct_ejecutado_validos[0], 1) if len(pct_ejecutado_validos) >= 2 else None

    return {
        'trimestre': trimestre, 'meses_fsc': meses_fsc, 'meses_planer': meses_planer,
        'variacion_pct_dentro': variacion_pct_dentro, 'variacion_pct_ejecutado': variacion_pct_ejecutado,
    }


def _combinar_subdirecciones(jerarquia_fsc, jerarquia_planer, fichas, detalle_fsc_subdireccion=None):
    """Combina el árbol de Jerarquía FSC (Dentro/Fuera + cumplimiento temporal) con
    el árbol de Ejecución PAC (Fichas) por NOMBRE de subdirección, para que el
    informe tenga UN capítulo por subdirección con ambas vistas — decisión del
    usuario 2026-07-21 ('un solo reporte unificado, capítulos por subdirección').
    Las 2 fuentes son deliberadamente independientes en el resto del sistema (ver
    docstrings de `calcular_pac_jerarquia`/`calcular_pac_jerarquia_planer`) — acá
    solo se juntan para efectos de presentación del informe, no se cruzan datos.

    `formularios_detalle`/`fichas_detalle`: listas planas (no agregadas por depto)
    para el detalle a nivel de FSC/proyecto individual dentro de cada capítulo —
    pedido 2026-07-23 para que el informe diga con claridad CON QUÉ formularios
    concretos se compone cada cifra, no solo el rollup por departamento."""
    detalle_fsc_subdireccion = detalle_fsc_subdireccion or {}
    fsc_por_nombre = {s['nombre']: s for s in jerarquia_fsc['subdirecciones']}
    planer_por_nombre = {s['nombre']: s for s in jerarquia_planer['subdirecciones']}
    nombres = list(dict.fromkeys(list(fsc_por_nombre.keys()) + list(planer_por_nombre.keys())))

    responsables_por_sub = {}
    fichas_por_sub = defaultdict(list)
    for f in fichas:
        if f['nombre_responsable'] and f['subdireccion_nombre']:
            responsables_por_sub.setdefault(f['subdireccion_nombre'], set()).add(f['nombre_responsable'])
        fichas_por_sub[f['subdireccion_nombre'] or 'Sin Clasificar'].append(f)

    combinado = [{
        'nombre': nombre,
        'fsc': fsc_por_nombre.get(nombre),
        'planer': planer_por_nombre.get(nombre),
        'responsables': sorted(responsables_por_sub.get(nombre, [])),
        'institucional': nombre in NOMBRES_SUBDIRECCIONES_INSTITUCIONALES,
        'formularios_detalle': detalle_fsc_subdireccion.get(nombre, []),
        'fichas_detalle': fichas_por_sub.get(nombre, []),
    } for nombre in nombres]

    combinado.sort(key=lambda s: (s['nombre'] == 'Sin Clasificar', not s['institucional'], s['nombre']))
    return combinado


# =============================================================================
# Word — helpers institucionales (portada, índice, pie de página, márgenes)
# =============================================================================

def _configurar_margenes_docx(doc):
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.3)
        section.left_margin = Cm(2.8)
        section.right_margin = Cm(2.2)


def _agregar_campo_docx(paragraph, instruccion):
    """Inserta un campo de Word de bajo nivel (PAGE, NUMPAGES, TOC) vía XML — la API
    de alto nivel de python-docx no soporta campos, solo texto/estilo."""
    run = paragraph.add_run()
    r = run._r
    fld_begin = OxmlElement('w:fldChar')
    fld_begin.set(qn('w:fldCharType'), 'begin')
    instr = OxmlElement('w:instrText')
    instr.set(qn('xml:space'), 'preserve')
    instr.text = instruccion
    fld_end = OxmlElement('w:fldChar')
    fld_end.set(qn('w:fldCharType'), 'end')
    r.append(fld_begin)
    r.append(instr)
    r.append(fld_end)


def _forzar_actualizacion_campos_docx(doc):
    """Sin esto, Word muestra el índice y los números de página con el texto
    placeholder hasta que el usuario actualiza los campos manualmente (F9). Con
    `w:updateFields`, Word los recalcula automáticamente al abrir el documento."""
    settings = doc.settings.element
    campo = OxmlElement('w:updateFields')
    campo.set(qn('w:val'), 'true')
    settings.append(campo)


def _agregar_pie_pagina_docx(doc):
    section = doc.sections[0]
    footer = section.footer
    footer.is_linked_to_previous = True
    p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    color_pie = RGBColor(0x94, 0xa3, 0xb8)

    run1 = p.add_run(f'{NOMBRE_INSTITUCION} · Informe de Seguimiento y Ejecución del Plan Anual de Compras · Página ')
    run1.font.size = Pt(8)
    run1.font.color.rgb = color_pie
    _agregar_campo_docx(p, 'PAGE')
    run2 = p.add_run(' de ')
    run2.font.size = Pt(8)
    run2.font.color.rgb = color_pie
    _agregar_campo_docx(p, 'NUMPAGES')


def _agregar_portada_docx(doc, periodo_label, anho, hoy):
    logo = _logo_bytes()
    if logo:
        p_logo = doc.add_paragraph()
        p_logo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_logo.add_run().add_picture(logo, width=Inches(1.5))

    doc.add_paragraph()
    titulo = doc.add_paragraph()
    titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_t = titulo.add_run(TITULO_INFORME)
    run_t.font.size = Pt(24)
    run_t.font.bold = True
    run_t.font.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_s = sub.add_run(NOMBRE_INSTITUCION.upper())
    run_s.font.size = Pt(15)
    run_s.font.bold = True
    run_s.font.color.rgb = RGBColor(0x38, 0xb2, 0xbd)

    edificio = _edificio_bytes()
    if edificio:
        doc.add_paragraph()
        p_img = doc.add_paragraph()
        p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_img.add_run().add_picture(edificio, width=Inches(5.6))

    doc.add_paragraph()
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_m = meta.add_run(
        f'Período de análisis: {periodo_label}\n'
        f'Año PAC (Plan Anual de Compras): {anho}\n'
        f'Generado el {hoy.strftime("%d-%m-%Y")}'
    )
    run_m.font.size = Pt(11.5)
    run_m.font.color.rgb = RGBColor(0x47, 0x55, 0x69)
    doc.add_page_break()


def _agregar_indice_docx(doc):
    h = doc.add_heading('Índice', level=1)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)
    p_ayuda = doc.add_paragraph()
    run_ayuda = p_ayuda.add_run(
        '(el índice se actualiza automáticamente al abrir el documento; si no se '
        've, haga clic derecho sobre él y seleccione "Actualizar campos")'
    )
    run_ayuda.font.size = Pt(9)
    run_ayuda.italic = True
    run_ayuda.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
    p_toc = doc.add_paragraph()
    _agregar_campo_docx(p_toc, 'TOC \\o "1-2" \\h \\z \\u')
    doc.add_page_break()


def _titulo_capitulo_docx(doc, texto):
    h = doc.add_heading(texto, level=1)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)
    return h


def _titulo_seccion_docx(doc, texto):
    h = doc.add_heading(texto, level=2)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0x38, 0x6f, 0xc9)
    return h


def _agregar_grafico_docx(doc, img, width, texto_caption):
    """Inserta un gráfico centrado con su leyenda (también centrada) debajo.
    `Document.add_picture()` crea un párrafo nuevo pero NO lo centra — a diferencia
    del patrón ya usado en la portada (`_agregar_portada_docx`), dejando los ~10
    gráficos del informe pegados al margen izquierdo en vez de centrados. Se
    centraliza acá para que todos los gráficos del documento luzcan parejos."""
    p_img = doc.add_paragraph()
    p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_img.add_run().add_picture(img, width=width)
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cap.add_run(texto_caption)
    run.italic = True
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)


# =============================================================================
# Word — tablas
# =============================================================================

def _tabla_departamentos(doc, departamentos):
    tabla = doc.add_table(rows=1, cols=5)
    tabla.style = 'Light Grid Accent 1'
    encabezados = ['Departamento', 'Total', '% Dentro PAC', '% En fecha', 'Monto Dentro PAC']
    for i, h in enumerate(encabezados):
        tabla.rows[0].cells[i].text = h
    for d in departamentos:
        fila = tabla.add_row().cells
        fila[0].text = _truncar(d['nombre'].title(), 55)
        fila[1].text = str(d['total'])
        fila[2].text = f"{d['pct_dentro']}%"
        fila[3].text = f"{d['pct_en_fecha']}%" if d['pct_en_fecha'] is not None else '—'
        fila[4].text = _money(d['monto_dentro'])


def _tabla_deptos_ejecucion(doc, departamentos):
    tabla = doc.add_table(rows=1, cols=5)
    tabla.style = 'Light Grid Accent 1'
    for i, h in enumerate(['Departamento', 'Total Fichas', 'Ejecutado', 'Pendiente', 'Atrasado']):
        tabla.rows[0].cells[i].text = h
    for dpt in departamentos:
        fila = tabla.add_row().cells
        fila[0].text = _truncar(dpt['nombre'].title(), 55)
        fila[1].text = str(dpt['total'])
        fila[2].text = str(dpt['ejecutados'])
        fila[3].text = str(dpt['pendientes'])
        fila[4].text = str(dpt['atrasados'])


def _tabla_ranking_depto(doc, filas):
    tabla = doc.add_table(rows=1, cols=5)
    tabla.style = 'Light Grid Accent 1'
    for i, h in enumerate(['Departamento', 'Subdirección', 'Total', 'Score', '% Dentro PAC']):
        tabla.rows[0].cells[i].text = h
    for f in filas:
        fila = tabla.add_row().cells
        fila[0].text = _truncar(f['nombre'].title(), 45)
        fila[1].text = _truncar(_nombre_subdireccion_display(f['subdireccion']), 35)
        fila[2].text = str(f['total'])
        fila[3].text = f"{f['score']:.0f}"
        fila[4].text = f"{f['pct_dentro']}%"


def _tabla_ranking_formulario(doc, filas):
    tabla = doc.add_table(rows=1, cols=5)
    tabla.style = 'Light Grid Accent 1'
    for i, h in enumerate(['Folio/Año', 'Unidad Requirente', 'Dentro/Fuera', 'Monto', 'Score']):
        tabla.rows[0].cells[i].text = h
    for f in filas:
        fila = tabla.add_row().cells
        fila[0].text = f"{f['folio']}/{f['anho']}"
        fila[1].text = _truncar(f['unidad_requirente'], 45)
        fila[2].text = 'Dentro' if f['dentro_fuera_pac'] == 'DENTRO' else 'Fuera'
        fila[3].text = _money(f['monto_estimado'])
        fila[4].text = f"{f['score']:.0f}"


def _tabla_resumen_subdireccion(doc, filas, anho, anho_anterior):
    tabla = doc.add_table(rows=1, cols=6)
    tabla.style = 'Light Grid Accent 1'
    for i, h in enumerate(['Subdirección', 'Dentro', 'Fuera', f'% Dentro {anho}', f'% Dentro {anho_anterior}', 'Variación']):
        tabla.rows[0].cells[i].text = h
    for f in filas:
        fila = tabla.add_row().cells
        fila[0].text = _truncar(_nombre_subdireccion_display(f['nombre']), 45)
        fila[1].text = str(f['dentro'])
        fila[2].text = str(f['fuera'])
        fila[3].text = f"{f['pct_dentro']}%"
        fila[4].text = f"{f['pct_dentro_anho_anterior']}%" if f['pct_dentro_anho_anterior'] is not None else '—'
        fila[5].text = f"{f['variacion_pp']:+.1f} pp" if f['variacion_pp'] is not None else '—'


def _tabla_fichas_pac_docx(doc, filas, limite=40):
    tabla = doc.add_table(rows=1, cols=6)
    tabla.style = 'Light Grid Accent 1'
    for i, h in enumerate(['ID Proyecto', 'Proyecto', 'Departamento', 'Responsable', 'Fecha Compra', 'Estado']):
        tabla.rows[0].cells[i].text = h
    for f in filas[:limite]:
        fila = tabla.add_row().cells
        fila[0].text = f['id_proyecto']
        fila[1].text = _truncar(f.get('nombre_proyecto') or '—', 45)
        fila[2].text = _truncar(f['depto_nombre'] or f['depto_texto'], 38)
        fila[3].text = _truncar(f['nombre_responsable'], 26)
        fila[4].text = f['fecha_mas_proxima'] or '—'
        fila[5].text = f['estado_ejecucion']
    if len(filas) > limite:
        p = doc.add_paragraph()
        run = p.add_run(f'… y {len(filas) - limite} ficha(s) adicional(es) — ver detalle completo en el dashboard interactivo.')
        run.italic = True
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)


def _tabla_proyectos_sin_iniciar_docx(doc, proyectos, limite=30):
    tabla = doc.add_table(rows=1, cols=4)
    tabla.style = 'Light Grid Accent 1'
    for i, h in enumerate(['ID Proyecto', 'Proyecto', 'Fecha Planificada', 'Estado']):
        tabla.rows[0].cells[i].text = h
    for p in proyectos[:limite]:
        fila = tabla.add_row().cells
        fila[0].text = p['id_proyecto']
        fila[1].text = _truncar(p.get('nombre_proyecto') or '—', 50)
        fila[2].text = p['fecha_inicio_compra']
        fila[3].text = 'Pendiente' if p['estado'] == 'PENDIENTE' else 'Atrasado'
    if len(proyectos) > limite:
        p_nota = doc.add_paragraph()
        run = p_nota.add_run(f'… y {len(proyectos) - limite} proyecto(s) adicional(es).')
        run.italic = True
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)


def _tabla_formularios_detalle_docx(doc, formularios, limite=30):
    """Detalle a nivel de CADA FSC individual (folio, unidad, monto) dentro del
    capítulo de una subdirección — complementa la tabla agregada por departamento
    para que el informe muestre con qué formularios concretos se compone la cifra."""
    tabla = doc.add_table(rows=1, cols=5)
    tabla.style = 'Light Grid Accent 1'
    for i, h in enumerate(['Formulario', 'Unidad Requirente', 'Dentro/Fuera', 'Monto', 'Fecha Derivado']):
        tabla.rows[0].cells[i].text = h
    for f in formularios[:limite]:
        fila = tabla.add_row().cells
        fila[0].text = f['id_formulario'] or f"{f['folio']}/{f['anho']}"
        fila[1].text = _truncar(f['unidad_requirente'], 45)
        fila[2].text = 'Dentro' if f['dentro_fuera_pac'] == 'DENTRO' else 'Fuera'
        fila[3].text = _money(f['monto_estimado'])
        fila[4].text = f['fecha_derivado'] or '—'
    if len(formularios) > limite:
        p = doc.add_paragraph()
        run = p.add_run(f'… y {len(formularios) - limite} formulario(s) adicional(es) — ver detalle completo en el dashboard interactivo.')
        run.italic = True
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)


def _tabla_fichas_ejecucion_detalle_docx(doc, fichas, limite=30):
    """Detalle a nivel de CADA proyecto/ficha del Plan de Compras dentro del capítulo
    de una subdirección, incluyendo CON QUÉ formulario(s) concreto(s) se ejecutó —
    pedido explícito 2026-07-23 ('hablar con claridad con qué FSC se ejecutaron
    los proyectos')."""
    tabla = doc.add_table(rows=1, cols=5)
    tabla.style = 'Light Grid Accent 1'
    for i, h in enumerate(['ID Proyecto', 'Proyecto', 'Departamento', 'Estado', 'Ejecutado con (Formulario)']):
        tabla.rows[0].cells[i].text = h
    for f in fichas[:limite]:
        fila = tabla.add_row().cells
        fila[0].text = f['id_proyecto']
        fila[1].text = _truncar(f['nombre_proyecto'] or '—', 42)
        fila[2].text = _truncar(f['depto_nombre'] or f['depto_texto'], 34)
        fila[3].text = f['estado_ejecucion']
        ejecutores = f.get('formularios_ejecutores') or []
        fila[4].text = ', '.join(fx['id_formulario'] or f"{fx['folio']}/{fx['anho']}" for fx in ejecutores) if ejecutores else '—'
    if len(fichas) > limite:
        p = doc.add_paragraph()
        run = p.add_run(f'… y {len(fichas) - limite} proyecto(s) adicional(es) — ver detalle completo en el dashboard interactivo.')
        run.italic = True
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)


def _render_grupos_subdireccion_docx(doc, grupos, tabla_fn, mensaje_vacio):
    """Renderiza una lista de grupos `{'nombre_display', 'items'}` (ver
    `_agrupar_por_subdireccion`) con un sub-encabezado por subdirección y su propia
    tabla — usado por Alertas y Seguimiento (2026-07-23) para que cada lista se lea
    por subdirección en vez de una sola tabla institucional mezclada."""
    if not grupos:
        doc.add_paragraph(mensaje_vacio)
        return
    for grupo in grupos:
        p = doc.add_paragraph()
        run = p.add_run(f'{grupo["nombre_display"]} ({_n(len(grupo["items"]))})')
        run.bold = True
        run.font.size = Pt(10.5)
        run.font.color.rgb = RGBColor(0x38, 0x6f, 0xc9)
        tabla_fn(doc, grupo['items'])


_PARRAFOS_METODOLOGIA = [
    'El presente informe integra dos fuentes de datos complementarias pero '
    'independientes entre sí, cada una respondiendo una pregunta distinta sobre '
    'la gestión del Plan Anual de Compras (PAC):',
    '(1) Cumplimiento Dentro/Fuera del PAC y Cumplimiento Temporal: se determina '
    'a nivel de cada Formulario de Solicitud de Compra (FSC) derivado a comprador, '
    'verificando si el proyecto que declara (ID de Plan) existe en el maestro '
    'histórico del Plan Anual de Compras. El cumplimiento temporal compara la '
    'fecha en que el formulario fue derivado contra la fecha de compra '
    'planificada para ese proyecto, con una tolerancia de un mes calendario.',
    '(2) Ejecución del Plan de Compras (Ficha PAC): se determina a nivel de cada '
    'ficha/proyecto del Plan Anual de Compras. Una ficha se considera "Ejecutada" '
    'cuando cuenta con al menos un formulario de compra u orden de compra '
    'enlazada; "Pendiente" cuando su fecha de compra planificada aún no vence; y '
    '"Atrasada" cuando la fecha ya venció sin formulario ni orden de compra '
    'asociada.',
    'Por tratarse de universos y momentos de corte distintos, las cifras de '
    'ambas fuentes no son directamente sumables entre sí — se presentan en '
    'secciones separadas dentro de cada capítulo de subdirección para mantener '
    'la trazabilidad de cada una.',
]


def _tabla_avance_trimestral_docx(doc, avance):
    """Tabla mes a mes del trimestre — % Dentro PAC y % Ejecutado del Plan de Compras,
    para visualizar la mejora dentro del propio trimestre (pedido 2026-07-22)."""
    meses_fsc_por_mes = {m['mes']: m for m in avance['meses_fsc']}
    meses_planer_por_mes = {m['mes']: m for m in avance['meses_planer']}
    meses_ids = sorted(set(meses_fsc_por_mes) | set(meses_planer_por_mes))

    tabla = doc.add_table(rows=1, cols=5)
    tabla.style = 'Light Grid Accent 1'
    for i, h in enumerate(['Mes', 'Formularios', '% Dentro PAC', 'Fichas Plan Compras', '% Ejecutado']):
        tabla.rows[0].cells[i].text = h
    for mes_id in meses_ids:
        mf, mp = meses_fsc_por_mes.get(mes_id), meses_planer_por_mes.get(mes_id)
        fila = tabla.add_row().cells
        fila[0].text = (mf or mp)['nombre_mes']
        fila[1].text = str(mf['total']) if mf else '—'
        fila[2].text = f"{mf['pct_dentro']}%" if mf and mf['total'] else '—'
        fila[3].text = str(mp['total']) if mp else '—'
        fila[4].text = f"{mp['pct_ejecutado']}%" if mp and mp['total'] else '—'


def _mejor_peor_de_ranking(rankings_depto, nombre_sub):
    mejor = next((f for f in rankings_depto['mejores'] if f['subdireccion'] == nombre_sub), None)
    peor = next((f for f in rankings_depto['peores'] if f['subdireccion'] == nombre_sub), None)
    return mejor, peor


def generar_informe_word(periodo):
    """Genera el informe Word institucional único y exhaustivo del período
    ('YYYY-MM' o 'YYYY-QN'). Retorna BytesIO. Estructura: Portada → Índice →
    Resumen Ejecutivo Institucional → Metodología → capítulos por subdirección
    (Formularios Dentro/Fuera PAC + Ejecución del Plan de Compras combinados) →
    Rankings Institucionales (departamentos y formularios) → Alertas y
    Seguimiento → Conclusiones → Anexo (otras unidades de la red)."""
    d = _datos_informe_completo(periodo)
    label, anho, anho_ant, hoy = d['label'], d['anho'], d['anho_anterior'], d['hoy']

    doc = Document()
    _configurar_margenes_docx(doc)

    # --- Portada + Índice ---------------------------------------------------
    _agregar_portada_docx(doc, label, anho, hoy)
    _agregar_indice_docx(doc)

    # --- Resumen Ejecutivo Institucional -------------------------------------
    _titulo_capitulo_docx(doc, 'Resumen Ejecutivo Institucional')
    doc.add_paragraph(parrafo_resumen_ejecutivo(label, d['dentro_fuera']['kpis'], d['comparativa_periodos']))
    doc.add_paragraph(parrafo_cumplimiento_temporal(label, d['temporal']['kpis']))
    doc.add_paragraph(
        f'En paralelo, el Plan Anual de Compras {anho} registra {_n(d["total_fichas"])} fichas/proyectos, '
        f'de las cuales {_n(d["ejecutados_fichas"])} ({d["pct_ejecutado_fichas"]}%) cuentan con formulario '
        f'de compra u orden de compra enlazada, por un monto total planificado de {_money(d["monto_total_fichas"])}.'
    )

    img = grafico_donut_dentro_fuera(d['dentro_fuera']['kpis']['pct_dentro'])
    if img:
        _agregar_grafico_docx(doc, img, Inches(2.8), 'Gráfico 1 — Distribución de formularios Dentro/Fuera del PAC en el período.')

    img = grafico_barras_comparativa_anual(d['dentro_fuera']['comparativa_anual'])
    if img:
        _agregar_grafico_docx(doc, img, Inches(5.6), 'Gráfico 2 — Comparativa histórica anual de formularios Dentro/Fuera del PAC.')

    img = grafico_evolucion_mensual_dentro_fuera(d['temporalidad_mensual']['meses'], anho, anho_ant)
    if img:
        _agregar_grafico_docx(
            doc, img, Inches(5.8),
            f'Gráfico 3 — Evolución mensual {anho} (barras) comparada contra el % Dentro del PAC de {anho_ant} '
            '(línea punteada), para visualizar la tendencia mes a mes respecto al año anterior.',
        )

    if d['avance_trimestral']:
        av = d['avance_trimestral']
        _titulo_seccion_docx(doc, f'Avance Trimestral — {label}')
        doc.add_paragraph(
            f'Detalle mes a mes de los 3 meses que componen el {label}, para visualizar la mejora o '
            'retroceso dentro del propio trimestre — complementario a la comparación anual del Gráfico 3.'
        )
        img = grafico_evolucion_mensual_dentro_fuera(av['meses_fsc'], anho, anho_ant, titulo=f'Avance mensual — {label}')
        if img:
            texto_var = (
                f'variación de {av["variacion_pct_dentro"]:+.1f} p.p. entre el primer y el último mes del trimestre'
                if av['variacion_pct_dentro'] is not None else 'sin datos suficientes para calcular variación interna'
            )
            _agregar_grafico_docx(doc, img, Inches(5.6), f'Gráfico 3b — % Dentro PAC mes a mes dentro del {label} ({texto_var}).')
        img = grafico_barras_ejecucion_mensual_planer(av['meses_planer'], titulo=f'Ejecución Plan de Compras — {label}')
        if img:
            texto_var = (
                f'variación de {av["variacion_pct_ejecutado"]:+.1f} p.p. entre el primer y el último mes'
                if av['variacion_pct_ejecutado'] is not None else 'sin datos suficientes para calcular variación interna'
            )
            _agregar_grafico_docx(doc, img, Inches(5.6), f'Gráfico 3c — % Ejecutado del Plan de Compras mes a mes dentro del {label} ({texto_var}).')
        _tabla_avance_trimestral_docx(doc, av)

    doc.add_page_break()

    _titulo_seccion_docx(doc, 'Cumplimiento Dentro/Fuera del PAC por Subdirección')
    doc.add_paragraph(
        'La siguiente tabla desglosa, por subdirección, el número de formularios Dentro y Fuera del PAC '
        f'del año {anho} comparado contra el mismo corte de {anho_ant}, para dimensionar la evolución del '
        'apego institucional al Plan Anual de Compras.'
    )
    if d['resumen_subdireccion']['subdirecciones']:
        _tabla_resumen_subdireccion(doc, d['resumen_subdireccion']['subdirecciones'], anho, anho_ant)
    img = grafico_donut_cumplimiento_temporal(d['temporal']['kpis'])
    if img:
        _agregar_grafico_docx(doc, img, Inches(3.4), 'Gráfico 4 — Cumplimiento temporal: formularios Dentro del PAC evaluados contra su fecha planificada.')

    img = grafico_barras_ejecucion_mensual_planer(d['temporal_mensual_planer']['meses'])
    if img:
        _agregar_grafico_docx(doc, img, Inches(5.8), f'Gráfico 5 — Ejecución mensual del Plan de Compras {anho} (fichas Ejecutadas/Pendientes/Atrasadas).')
    doc.add_page_break()

    # --- Metodología ----------------------------------------------------------
    _titulo_capitulo_docx(doc, 'Metodología')
    for parrafo in _PARRAFOS_METODOLOGIA:
        doc.add_paragraph(parrafo)
    doc.add_page_break()

    # --- Capítulos por Subdirección --------------------------------------------
    subs_institucionales = [s for s in d['subdirecciones'] if s['institucional']]
    otras_ramas = [s for s in d['subdirecciones'] if not s['institucional']]

    for sub in subs_institucionales:
        nombre_display = _nombre_subdireccion_display(sub['nombre'])
        _titulo_capitulo_docx(doc, nombre_display)

        _titulo_seccion_docx(doc, f'Formularios de Solicitud de Compra — Período {label}')
        if sub['fsc'] and sub['fsc']['total']:
            mejor_rk, peor_rk = _mejor_peor_de_ranking(d['rankings_depto'], sub['nombre'])
            doc.add_paragraph(parrafo_capitulo_subdireccion(nombre_display, sub['fsc'], mejor_rk, peor_rk))
            if sub['fsc']['departamentos']:
                _tabla_departamentos(doc, sub['fsc']['departamentos'])

            img = grafico_donut_dentro_fuera(sub['fsc']['pct_dentro'], titulo=f'{nombre_display} — Dentro/Fuera PAC')
            if img:
                _agregar_grafico_docx(
                    doc, img, Inches(2.6),
                    f'Distribución de los {_n(sub["fsc"]["total"])} formularios de {nombre_display} entre Dentro '
                    'y Fuera del PAC durante el período.',
                )

            if sub['formularios_detalle']:
                p_intro = doc.add_paragraph()
                run = p_intro.add_run(
                    f'Detalle de los {_n(len(sub["formularios_detalle"]))} formularios individuales de '
                    f'{nombre_display} en el período:'
                )
                run.font.size = Pt(9.5)
                _tabla_formularios_detalle_docx(doc, sub['formularios_detalle'])
        else:
            doc.add_paragraph('Sin formularios registrados en el período para esta subdirección.')

        _titulo_seccion_docx(doc, f'Ejecución del Plan de Compras — Año PAC {anho}')
        if sub['planer'] and sub['planer']['total']:
            p = sub['planer']
            doc.add_paragraph(
                f"{_n(p['total'])} fichas del Plan de Compras — {p['pct_ejecutado']}% ejecutadas "
                f"({_n(p['ejecutados'])} ejecutadas, {_n(p['pendientes'])} pendientes, {_n(p['atrasados'])} atrasadas). "
                f"Monto total planificado: {_money(p['monto_total'])}."
            )
            if p['departamentos']:
                _tabla_deptos_ejecucion(doc, p['departamentos'])
            if sub['responsables']:
                p_resp = doc.add_paragraph()
                run = p_resp.add_run('Responsables: ' + ', '.join(sub['responsables']))
                run.font.size = Pt(9.5)
                run.font.color.rgb = RGBColor(0x47, 0x55, 0x69)

            img = grafico_donut_ejecucion_planer(
                p['ejecutados'], p['pendientes'], p['atrasados'], titulo=f'{nombre_display} — Ejecución del Plan de Compras',
            )
            if img:
                _agregar_grafico_docx(
                    doc, img, Inches(3.0),
                    f'Estado de ejecución de los {_n(p["total"])} proyectos del Plan de Compras {anho} de {nombre_display}.',
                )

            if sub['fichas_detalle']:
                p_intro = doc.add_paragraph()
                run = p_intro.add_run(
                    f'Detalle de los {_n(len(sub["fichas_detalle"]))} proyectos del Plan de Compras {anho} de '
                    f'{nombre_display}, indicando con qué formulario(s) se ejecutó cada uno:'
                )
                run.font.size = Pt(9.5)
                _tabla_fichas_ejecucion_detalle_docx(doc, sub['fichas_detalle'])
        else:
            doc.add_paragraph('Sin fichas del Plan de Compras registradas para esta subdirección.')
        doc.add_page_break()

    # --- Rankings Institucionales -----------------------------------------------
    _titulo_capitulo_docx(doc, 'Rankings Institucionales')
    doc.add_paragraph(
        'El score compuesto de cada ranking pondera 40% el % Dentro del PAC, 40% el % de cumplimiento '
        'temporal (fecha de derivación vs. fecha planificada) y 20% el % Dentro del PAC ponderado por monto '
        '— evitando que el resultado dependa solo del tamaño del presupuesto gestionado.'
    )
    _titulo_seccion_docx(doc, 'Mejores departamentos')
    if d['rankings_depto']['mejores']:
        _tabla_ranking_depto(doc, d['rankings_depto']['mejores'])
    else:
        doc.add_paragraph('Sin datos suficientes para generar el ranking en este período.')
    _titulo_seccion_docx(doc, 'Departamentos con mayor oportunidad de mejora')
    if d['rankings_depto']['peores']:
        _tabla_ranking_depto(doc, d['rankings_depto']['peores'])
    else:
        doc.add_paragraph('Sin datos suficientes para generar el ranking en este período.')

    _titulo_seccion_docx(doc, 'Mejores formularios')
    if d['rankings_formulario']['mejores']:
        _tabla_ranking_formulario(doc, d['rankings_formulario']['mejores'])
    else:
        doc.add_paragraph('Sin datos suficientes para generar el ranking en este período.')
    _titulo_seccion_docx(doc, 'Formularios con mayor oportunidad de mejora')
    if d['rankings_formulario']['peores']:
        _tabla_ranking_formulario(doc, d['rankings_formulario']['peores'])
    else:
        doc.add_paragraph('Sin datos suficientes para generar el ranking en este período.')
    doc.add_page_break()

    # --- Alertas y Seguimiento ----------------------------------------------
    _titulo_capitulo_docx(doc, 'Alertas y Seguimiento')
    doc.add_paragraph(
        'Las secciones siguientes se agrupan por subdirección institucional para facilitar el '
        'seguimiento y la coordinación con cada responsable.'
    )
    _titulo_seccion_docx(doc, 'Proyectos planificados sin ningún formulario Dentro PAC derivado todavía')
    _render_grupos_subdireccion_docx(
        doc, d['proyectos_sin_iniciar'], _tabla_proyectos_sin_iniciar_docx,
        'No hay proyectos planificados sin iniciar en el período.',
    )

    _titulo_seccion_docx(doc, f'Fichas a ejecutar el próximo mes (Año PAC {anho})')
    _render_grupos_subdireccion_docx(
        doc, d['proximo_mes'], _tabla_fichas_pac_docx,
        'No hay fichas planificadas para el próximo mes.',
    )

    _titulo_seccion_docx(doc, f'Fichas atrasadas (Año PAC {anho})')
    _render_grupos_subdireccion_docx(
        doc, d['atrasadas'], _tabla_fichas_pac_docx,
        'No hay fichas atrasadas en el período.',
    )
    doc.add_page_break()

    # --- Conclusiones y Recomendaciones ---------------------------------------
    _titulo_capitulo_docx(doc, 'Conclusiones y Recomendaciones')
    doc.add_paragraph(parrafo_conclusiones(label, d['dentro_fuera']['kpis'], d['temporal']['kpis']))
    doc.add_paragraph(
        f'En cuanto a la ejecución del Plan de Compras {anho}, {d["pct_ejecutado_fichas"]}% de las fichas '
        f'planificadas cuentan hoy con formulario u orden de compra enlazada. '
        + ('Se recomienda priorizar la revisión de las fichas atrasadas listadas en la sección de Alertas '
           'y coordinar con los responsables de cada subdirección la regularización de los proyectos pendientes.'
           if d['pct_ejecutado_fichas'] < 70 else
           'El nivel de ejecución es satisfactorio; se recomienda mantener el seguimiento mensual para '
           'sostener la tendencia.')
    )
    doc.add_page_break()

    # --- Anexo: otras unidades de la red ---------------------------------------
    _titulo_capitulo_docx(doc, 'Anexo — Otras Unidades de la Red')
    doc.add_paragraph(
        'Detalle abreviado de las restantes unidades de la red asistencial (fuera de la Dirección Servicio '
        'de Salud Osorno) y de los formularios que no pudieron clasificarse organizacionalmente.'
    )
    if otras_ramas:
        for sub in otras_ramas:
            _titulo_seccion_docx(doc, _nombre_subdireccion_display(sub['nombre']))
            if sub['fsc'] and sub['fsc']['departamentos']:
                _tabla_departamentos(doc, sub['fsc']['departamentos'])
            if sub['planer'] and sub['planer']['departamentos']:
                doc.add_paragraph('Ejecución del Plan de Compras:')
                _tabla_deptos_ejecucion(doc, sub['planer']['departamentos'])
    else:
        doc.add_paragraph('Sin registros de otras unidades de la red en este período.')

    _agregar_pie_pagina_docx(doc)
    _forzar_actualizacion_campos_docx(doc)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# =============================================================================
# PDF — reportlab (Platypus). Índice con numeración de página real vía
# `BaseDocTemplate.multiBuild()` (2 pasadas: la 1ª resuelve en qué página cae
# cada capítulo, la 2ª ya puede imprimir el índice con esos números).
# =============================================================================

_PDF_ESTILOS = getSampleStyleSheet()
_PDF_ESTILOS.add(ParagraphStyle(
    name='TituloPortada', fontSize=22, leading=28, alignment=1, spaceAfter=10,
    textColor=colors.HexColor(COLOR_INSTITUCIONAL), fontName='Helvetica-Bold',
))
_PDF_ESTILOS.add(ParagraphStyle(
    name='SubportadaFuerte', fontSize=14, leading=18, alignment=1, spaceAfter=6,
    textColor=colors.HexColor(COLOR_INSTITUCIONAL_CLARO), fontName='Helvetica-Bold',
))
_PDF_ESTILOS.add(ParagraphStyle(
    name='Subportada', fontSize=11.5, leading=17, alignment=1, textColor=colors.HexColor('#475569'),
))
_PDF_ESTILOS.add(ParagraphStyle(
    name='TituloCapitulo', fontSize=16, leading=20, spaceBefore=6, spaceAfter=10,
    textColor=colors.HexColor(COLOR_INSTITUCIONAL), fontName='Helvetica-Bold',
))
_PDF_ESTILOS.add(ParagraphStyle(
    name='TituloSeccion', fontSize=12.5, leading=16, spaceBefore=8, spaceAfter=8,
    textColor=colors.HexColor('#386fc9'), fontName='Helvetica-Bold',
))
_PDF_ESTILOS.add(ParagraphStyle(
    name='SubGrupo', fontSize=10, leading=13, spaceBefore=6, spaceAfter=3,
    textColor=colors.HexColor(COLOR_INSTITUCIONAL), fontName='Helvetica-Bold',
))
_PDF_ESTILOS.add(ParagraphStyle(
    name='Cuerpo', fontSize=10, leading=14.5, spaceAfter=9, textColor=colors.HexColor('#374151'),
))
_PDF_ESTILOS.add(ParagraphStyle(
    name='Leyenda', fontSize=8, leading=11, alignment=1, spaceAfter=10,
    textColor=colors.HexColor('#94a3b8'), fontName='Helvetica-Oblique',
))
_PDF_ESTILOS.add(ParagraphStyle(
    name='Metadato', fontSize=8, leading=11, spaceAfter=10, textColor=colors.HexColor('#475569'),
))
_PDF_ESTILOS.add(ParagraphStyle(
    name='Celda', fontSize=8, leading=10, textColor=colors.HexColor('#334155'),
))
_PDF_ESTILOS.add(ParagraphStyle(
    name='CeldaEncabezado', fontSize=8, leading=10, textColor=colors.HexColor('#334155'), fontName='Helvetica-Bold',
))

_PDF_TABLA_ESTILO = TableStyle([
    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#f1f5f9')),
    ('FONTSIZE', (0, 0), (-1, -1), 8),
    ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#e2e8f0')),
    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f8fafc')]),
    ('TOPPADDING', (0, 0), (-1, -1), 4),
    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
    ('LEFTPADDING', (0, 0), (-1, -1), 5),
    ('RIGHTPADDING', (0, 0), (-1, -1), 5),
])


def _celda(texto):
    """Envuelve el texto en un Paragraph para que reportlab lo ajuste (word-wrap)
    dentro del ancho de columna en vez de desbordarlo — pedido explícito del
    usuario ('ojo con las tablas que no se desborde el texto')."""
    return Paragraph(str(texto if texto not in (None, '') else '—'), _PDF_ESTILOS['Celda'])


def _celda_encabezado(texto):
    return Paragraph(str(texto), _PDF_ESTILOS['CeldaEncabezado'])


def _fila_encabezado(encabezados):
    return [_celda_encabezado(h) for h in encabezados]


def _pdf_imagen(img_bytes, width_in, height_in):
    if not img_bytes:
        return None
    img = RLImage(img_bytes, width=width_in * inch, height=height_in * inch)
    img.hAlign = 'CENTER'
    return img


def _pdf_tabla_departamentos(departamentos):
    """Sin `limite` — a diferencia de versiones anteriores (`limite=20`), que
    truncaban silenciosamente subdirecciones con más de 20 departamentos mientras
    su par Word (`_tabla_departamentos`) mostraba la lista completa (bug real,
    revisión de código 2026-07-27). Se iguala al comportamiento de Word: mostrar
    todos los departamentos, sin límite."""
    filas = [_fila_encabezado(['Departamento', 'Total', '% Dentro', '% En fecha', 'Monto Dentro PAC'])]
    for d in departamentos:
        filas.append([
            _celda(d['nombre'].title()), str(d['total']), f"{d['pct_dentro']}%",
            f"{d['pct_en_fecha']}%" if d['pct_en_fecha'] is not None else '—', _money(d['monto_dentro']),
        ])
    tabla = Table(filas, colWidths=[2.3 * inch, 0.55 * inch, 0.75 * inch, 0.75 * inch, 1.35 * inch], repeatRows=1)
    tabla.setStyle(_PDF_TABLA_ESTILO)
    return tabla


def _pdf_tabla_deptos_ejecucion(departamentos):
    """Sin `limite` — ver docstring de `_pdf_tabla_departamentos` (mismo bug,
    misma corrección: paridad con el Word sin límite)."""
    filas = [_fila_encabezado(['Departamento', 'Total Fichas', 'Ejecutado', 'Pendiente', 'Atrasado'])]
    for dpt in departamentos:
        filas.append([_celda(dpt['nombre'].title()), str(dpt['total']), str(dpt['ejecutados']), str(dpt['pendientes']), str(dpt['atrasados'])])
    tabla = Table(filas, colWidths=[2.7 * inch, 0.7 * inch, 0.75 * inch, 0.75 * inch, 0.75 * inch], repeatRows=1)
    tabla.setStyle(_PDF_TABLA_ESTILO)
    return tabla


def _pdf_tabla_ranking_depto(filas_ranking, limite=15):
    filas = [_fila_encabezado(['Departamento', 'Subdirección', 'Total', 'Score', '% Dentro'])]
    for f in filas_ranking[:limite]:
        filas.append([_celda(f['nombre'].title()), _celda(_nombre_subdireccion_display(f['subdireccion'])), str(f['total']), f"{f['score']:.0f}", f"{f['pct_dentro']}%"])
    tabla = Table(filas, colWidths=[2 * inch, 2 * inch, 0.6 * inch, 0.6 * inch, 0.7 * inch], repeatRows=1)
    tabla.setStyle(_PDF_TABLA_ESTILO)
    return tabla


def _pdf_tabla_ranking_formulario(filas_ranking, limite=15):
    filas = [_fila_encabezado(['Folio/Año', 'Unidad Requirente', 'Dentro/Fuera', 'Monto', 'Score'])]
    for f in filas_ranking[:limite]:
        filas.append([
            f"{f['folio']}/{f['anho']}", _celda(f['unidad_requirente']),
            'Dentro' if f['dentro_fuera_pac'] == 'DENTRO' else 'Fuera', _money(f['monto_estimado']), f"{f['score']:.0f}",
        ])
    tabla = Table(filas, colWidths=[0.8 * inch, 2.6 * inch, 0.85 * inch, 1 * inch, 0.55 * inch], repeatRows=1)
    tabla.setStyle(_PDF_TABLA_ESTILO)
    return tabla


def _pdf_tabla_resumen_subdireccion(filas, anho, anho_anterior, limite=15):
    encabezado = ['Subdirección', 'Dentro', 'Fuera', f'% Dentro {anho}', f'% Dentro {anho_anterior}', 'Variación']
    tabla_filas = [_fila_encabezado(encabezado)]
    for f in filas[:limite]:
        tabla_filas.append([
            _celda(_nombre_subdireccion_display(f['nombre'])), str(f['dentro']), str(f['fuera']), f"{f['pct_dentro']}%",
            f"{f['pct_dentro_anho_anterior']}%" if f['pct_dentro_anho_anterior'] is not None else '—',
            f"{f['variacion_pp']:+.1f} pp" if f['variacion_pp'] is not None else '—',
        ])
    tabla = Table(tabla_filas, colWidths=[2.2 * inch, 0.6 * inch, 0.6 * inch, 0.9 * inch, 0.9 * inch, 0.8 * inch], repeatRows=1)
    tabla.setStyle(_PDF_TABLA_ESTILO)
    return tabla


def _pdf_nota_adicionales(total, limite, texto_unidad):
    """Paragraph '… y N adicional(es)' cuando una tabla PDF trunca una lista más larga
    que `limite` — paridad con el aviso equivalente que ya llevan las tablas Word
    truncadas (`_tabla_fichas_pac_docx`, `_tabla_proyectos_sin_iniciar_docx`, etc.);
    el PDF no lo tenía en ninguna de sus tablas (bug real, revisión de código
    2026-07-27), dejando al lector sin forma de saber que hubo recorte. Devuelve
    `None` (nada que agregar al story) cuando no hubo truncamiento."""
    if total <= limite:
        return None
    return Paragraph(
        f'… y {total - limite} {texto_unidad} adicional(es) — ver detalle completo en el dashboard interactivo.',
        _PDF_ESTILOS['Leyenda'],
    )


def _pdf_tabla_fichas_pac(filas, limite=40):
    """`limite=40` — antes 25, mientras el Word equivalente (`_tabla_fichas_pac_docx`)
    usa 40 para la MISMA fuente de datos (`d['proximo_mes']`/`d['atrasadas']`); el PDF
    podía omitir hasta 15 fichas que el Word sí mostraba (bug real, revisión de código
    2026-07-27). Ahora retorna `[tabla, nota_o_none]` — ver `_pdf_nota_adicionales`."""
    encabezado = ['ID Proyecto', 'Proyecto', 'Departamento', 'Responsable', 'Fecha Compra', 'Estado']
    tabla_filas = [_fila_encabezado(encabezado)]
    for f in filas[:limite]:
        tabla_filas.append([
            f['id_proyecto'], _celda(f.get('nombre_proyecto') or '—'), _celda(f['depto_nombre'] or f['depto_texto']),
            _celda(f['nombre_responsable']), f['fecha_mas_proxima'] or '—', f['estado_ejecucion'],
        ])
    tabla = Table(tabla_filas, colWidths=[0.75 * inch, 1.55 * inch, 1.5 * inch, 1.15 * inch, 0.8 * inch, 0.7 * inch], repeatRows=1)
    tabla.setStyle(_PDF_TABLA_ESTILO)
    nota = _pdf_nota_adicionales(len(filas), limite, 'ficha(s)')
    return [tabla, nota] if nota else [tabla]


def _pdf_tabla_formularios_detalle(formularios, limite=30):
    """Detalle a nivel de CADA FSC individual — complemento reportlab de
    `_tabla_formularios_detalle_docx` para el capítulo de subdirección del PDF.
    Retorna `[tabla, nota_o_none]` — ver `_pdf_nota_adicionales`."""
    encabezado = ['Formulario', 'Unidad Requirente', 'Dentro/Fuera', 'Monto', 'Fecha Derivado']
    tabla_filas = [_fila_encabezado(encabezado)]
    for f in formularios[:limite]:
        tabla_filas.append([
            f['id_formulario'] or f"{f['folio']}/{f['anho']}", _celda(f['unidad_requirente']),
            'Dentro' if f['dentro_fuera_pac'] == 'DENTRO' else 'Fuera', _money(f['monto_estimado']),
            f['fecha_derivado'] or '—',
        ])
    tabla = Table(tabla_filas, colWidths=[0.95 * inch, 2.5 * inch, 0.85 * inch, 1.1 * inch, 1.05 * inch], repeatRows=1)
    tabla.setStyle(_PDF_TABLA_ESTILO)
    nota = _pdf_nota_adicionales(len(formularios), limite, 'formulario(s)')
    return [tabla, nota] if nota else [tabla]


def _pdf_tabla_fichas_ejecucion_detalle(fichas, limite=30):
    """Detalle de proyectos del Plan de Compras con CON QUÉ formulario(s) se
    ejecutó cada uno — complemento reportlab de `_tabla_fichas_ejecucion_detalle_docx`.
    Retorna `[tabla, nota_o_none]` — ver `_pdf_nota_adicionales`."""
    encabezado = ['ID Proyecto', 'Proyecto', 'Departamento', 'Estado', 'Ejecutado con (Formulario)']
    tabla_filas = [_fila_encabezado(encabezado)]
    for f in fichas[:limite]:
        ejecutores = f.get('formularios_ejecutores') or []
        texto_ejecutores = ', '.join(fx['id_formulario'] or f"{fx['folio']}/{fx['anho']}" for fx in ejecutores) if ejecutores else '—'
        tabla_filas.append([
            f['id_proyecto'], _celda(f['nombre_proyecto'] or '—'), _celda(f['depto_nombre'] or f['depto_texto']),
            f['estado_ejecucion'], _celda(texto_ejecutores),
        ])
    tabla = Table(tabla_filas, colWidths=[0.75 * inch, 1.5 * inch, 1.35 * inch, 0.75 * inch, 2.15 * inch], repeatRows=1)
    tabla.setStyle(_PDF_TABLA_ESTILO)
    nota = _pdf_nota_adicionales(len(fichas), limite, 'proyecto(s)')
    return [tabla, nota] if nota else [tabla]


def _pdf_tabla_avance_trimestral(avance):
    """Tabla mes a mes del trimestre — misma lógica que `_tabla_avance_trimestral_docx`
    (Word), en formato reportlab con celdas `Paragraph` (sin desborde de texto)."""
    meses_fsc_por_mes = {m['mes']: m for m in avance['meses_fsc']}
    meses_planer_por_mes = {m['mes']: m for m in avance['meses_planer']}
    meses_ids = sorted(set(meses_fsc_por_mes) | set(meses_planer_por_mes))

    tabla_filas = [_fila_encabezado(['Mes', 'Formularios', '% Dentro PAC', 'Fichas Plan Compras', '% Ejecutado'])]
    for mes_id in meses_ids:
        mf, mp = meses_fsc_por_mes.get(mes_id), meses_planer_por_mes.get(mes_id)
        tabla_filas.append([
            (mf or mp)['nombre_mes'],
            str(mf['total']) if mf else '—',
            f"{mf['pct_dentro']}%" if mf and mf['total'] else '—',
            str(mp['total']) if mp else '—',
            f"{mp['pct_ejecutado']}%" if mp and mp['total'] else '—',
        ])
    tabla = Table(tabla_filas, colWidths=[1 * inch, 1.1 * inch, 1.1 * inch, 1.3 * inch, 1.1 * inch], repeatRows=1)
    tabla.setStyle(_PDF_TABLA_ESTILO)
    return tabla


def _pdf_tabla_proyectos_sin_iniciar(proyectos, limite=30):
    """Retorna `[tabla, nota_o_none]` — ver `_pdf_nota_adicionales`."""
    tabla_filas = [_fila_encabezado(['ID Proyecto', 'Proyecto', 'Fecha Planificada', 'Estado'])]
    for p in proyectos[:limite]:
        tabla_filas.append([
            p['id_proyecto'], _celda(p.get('nombre_proyecto') or '—'),
            p['fecha_inicio_compra'], 'Pendiente' if p['estado'] == 'PENDIENTE' else 'Atrasado',
        ])
    tabla = Table(tabla_filas, colWidths=[1 * inch, 2.7 * inch, 1.3 * inch, 0.9 * inch], repeatRows=1)
    tabla.setStyle(_PDF_TABLA_ESTILO)
    nota = _pdf_nota_adicionales(len(proyectos), limite, 'proyecto(s)')
    return [tabla, nota] if nota else [tabla]


def _render_grupos_subdireccion_pdf(story, grupos, tabla_fn, est, mensaje_vacio):
    """Análogo reportlab de `_render_grupos_subdireccion_docx`: agrega un
    sub-encabezado por subdirección y su tabla al `story` — Alertas y Seguimiento
    (2026-07-23) agrupado por subdirección en vez de una sola tabla mezclada."""
    if not grupos:
        story.append(Paragraph(mensaje_vacio, est['Cuerpo']))
        return
    for grupo in grupos:
        story.append(Paragraph(f'{grupo["nombre_display"]} ({_n(len(grupo["items"]))})', est['SubGrupo']))
        story.extend(tabla_fn(grupo['items']))
        story.append(Spacer(1, 0.12 * inch))


class _InformeDocTemplate(BaseDocTemplate):
    """DocTemplate con pie de página institucional en todas las hojas y registro
    de entradas de índice (`notify('TOCEntry', ...)`) leyendo el estilo de cada
    Paragraph — 'TituloCapitulo' → nivel 0, 'TituloSeccion' → nivel 1. Requiere
    `multiBuild()` en vez de `build()` para que el índice resuelva números de
    página reales (reportlab necesita una pasada extra para saber en qué página
    terminó cada capítulo antes de poder imprimir el índice)."""

    def __init__(self, *args, **kwargs):
        BaseDocTemplate.__init__(self, *args, **kwargs)
        frame = Frame(self.leftMargin, self.bottomMargin, self.width, self.height, id='normal')
        self.addPageTemplates([PageTemplate(id='con_pie', frames=[frame], onPage=self._dibujar_pie)])

    def _dibujar_pie(self, canvas, doc):
        canvas.saveState()
        canvas.setFont('Helvetica', 7.5)
        canvas.setFillColor(colors.HexColor('#94a3b8'))
        canvas.drawCentredString(
            doc.pagesize[0] / 2, 0.45 * inch,
            f'{NOMBRE_INSTITUCION} · Informe de Seguimiento y Ejecución del Plan Anual de Compras · Página {doc.page}',
        )
        canvas.restoreState()

    def afterFlowable(self, flowable):
        if isinstance(flowable, Paragraph):
            estilo = flowable.style.name
            texto = flowable.getPlainText()
            if estilo == 'TituloCapitulo':
                self.notify('TOCEntry', (0, texto, self.page))
            elif estilo == 'TituloSeccion':
                self.notify('TOCEntry', (1, texto, self.page))


def _portada_pdf(story, periodo_label, anho, hoy):
    logo = _pdf_imagen(_logo_bytes(), 1.3, 1.3)
    if logo:
        story.append(logo)
    story.append(Spacer(1, 0.25 * inch))
    story.append(Paragraph(TITULO_INFORME, _PDF_ESTILOS['TituloPortada']))
    story.append(Paragraph(NOMBRE_INSTITUCION.upper(), _PDF_ESTILOS['SubportadaFuerte']))
    edificio = _pdf_imagen(_edificio_bytes(), 5.3, 3.48)
    if edificio:
        story.append(Spacer(1, 0.2 * inch))
        story.append(edificio)
    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph(
        f'Período de análisis: {periodo_label}<br/>Año PAC (Plan Anual de Compras): {anho}<br/>'
        f'Generado el {hoy.strftime("%d-%m-%Y")}',
        _PDF_ESTILOS['Subportada'],
    ))
    story.append(PageBreak())


def _indice_pdf(story):
    toc = TableOfContents()
    toc.levelStyles = [
        ParagraphStyle(name='TOCCapitulo', fontSize=11, leading=16, textColor=colors.HexColor(COLOR_INSTITUCIONAL), spaceBefore=4),
        ParagraphStyle(name='TOCSeccion', fontSize=9.5, leading=13, leftIndent=16, textColor=colors.HexColor('#475569')),
    ]
    story.append(Paragraph('Índice', _PDF_ESTILOS['TituloCapitulo']))
    story.append(toc)
    story.append(PageBreak())


def generar_reporte_pdf(periodo):
    """Genera el reporte PDF institucional único y exhaustivo del período
    ('YYYY-MM' o 'YYYY-QN'). Retorna BytesIO. Misma estructura de capítulos que
    `generar_informe_word` — ver ese docstring."""
    d = _datos_informe_completo(periodo)
    label, anho, anho_ant, hoy = d['label'], d['anho'], d['anho_anterior'], d['hoy']
    est = _PDF_ESTILOS

    buf = io.BytesIO()
    doc = _InformeDocTemplate(
        buf, pagesize=letter,
        topMargin=0.9 * inch, bottomMargin=0.85 * inch, leftMargin=0.95 * inch, rightMargin=0.75 * inch,
    )
    story = []

    _portada_pdf(story, label, anho, hoy)
    _indice_pdf(story)

    # --- Resumen Ejecutivo Institucional -------------------------------------
    story.append(Paragraph('Resumen Ejecutivo Institucional', est['TituloCapitulo']))
    story.append(Paragraph(parrafo_resumen_ejecutivo(label, d['dentro_fuera']['kpis'], d['comparativa_periodos']), est['Cuerpo']))
    story.append(Paragraph(parrafo_cumplimiento_temporal(label, d['temporal']['kpis']), est['Cuerpo']))
    story.append(Paragraph(
        f'En paralelo, el Plan Anual de Compras {anho} registra {_n(d["total_fichas"])} fichas/proyectos, de las '
        f'cuales {_n(d["ejecutados_fichas"])} ({d["pct_ejecutado_fichas"]}%) cuentan con formulario de compra u '
        f'orden de compra enlazada, por un monto total planificado de {_money(d["monto_total_fichas"])}.',
        est['Cuerpo'],
    ))

    img = _pdf_imagen(grafico_donut_dentro_fuera(d['dentro_fuera']['kpis']['pct_dentro']), 2.5, 2.5)
    if img:
        story.append(img)
        story.append(Paragraph('Gráfico 1 — Distribución de formularios Dentro/Fuera del PAC en el período.', est['Leyenda']))
    img = _pdf_imagen(grafico_barras_comparativa_anual(d['dentro_fuera']['comparativa_anual']), 5.4, 3.0)
    if img:
        story.append(img)
        story.append(Paragraph('Gráfico 2 — Comparativa histórica anual de formularios Dentro/Fuera del PAC.', est['Leyenda']))
    img = _pdf_imagen(grafico_evolucion_mensual_dentro_fuera(d['temporalidad_mensual']['meses'], anho, anho_ant), 5.6, 2.6)
    if img:
        story.append(img)
        story.append(Paragraph(
            f'Gráfico 3 — Evolución mensual {anho} (barras) comparada contra el % Dentro del PAC de {anho_ant} '
            '(línea punteada).', est['Leyenda'],
        ))

    if d['avance_trimestral']:
        av = d['avance_trimestral']
        story.append(Paragraph(f'Avance Trimestral — {label}', est['TituloSeccion']))
        story.append(Paragraph(
            f'Detalle mes a mes de los 3 meses que componen el {label}, para visualizar la mejora o retroceso '
            'dentro del propio trimestre — complementario a la comparación anual del Gráfico 3.', est['Cuerpo'],
        ))
        img = _pdf_imagen(grafico_evolucion_mensual_dentro_fuera(av['meses_fsc'], anho, anho_ant, titulo=f'Avance mensual — {label}'), 5.6, 2.6)
        if img:
            story.append(img)
            texto_var = (
                f'variación de {av["variacion_pct_dentro"]:+.1f} p.p. entre el primer y el último mes del trimestre'
                if av['variacion_pct_dentro'] is not None else 'sin datos suficientes para calcular variación interna'
            )
            story.append(Paragraph(f'Gráfico 3b — % Dentro PAC mes a mes dentro del {label} ({texto_var}).', est['Leyenda']))
        img = _pdf_imagen(grafico_barras_ejecucion_mensual_planer(av['meses_planer'], titulo=f'Ejecución Plan de Compras — {label}'), 5.6, 2.6)
        if img:
            story.append(img)
            texto_var = (
                f'variación de {av["variacion_pct_ejecutado"]:+.1f} p.p. entre el primer y el último mes'
                if av['variacion_pct_ejecutado'] is not None else 'sin datos suficientes para calcular variación interna'
            )
            story.append(Paragraph(f'Gráfico 3c — % Ejecutado del Plan de Compras mes a mes dentro del {label} ({texto_var}).', est['Leyenda']))
        story.append(_pdf_tabla_avance_trimestral(av))

    story.append(PageBreak())

    story.append(Paragraph('Cumplimiento Dentro/Fuera del PAC por Subdirección', est['TituloSeccion']))
    story.append(Paragraph(
        'Desglose por subdirección del número de formularios Dentro y Fuera del PAC del año en curso comparado '
        'contra el mismo corte del año anterior.', est['Cuerpo'],
    ))
    if d['resumen_subdireccion']['subdirecciones']:
        story.append(_pdf_tabla_resumen_subdireccion(d['resumen_subdireccion']['subdirecciones'], anho, anho_ant))
    img = _pdf_imagen(grafico_donut_cumplimiento_temporal(d['temporal']['kpis']), 3.0, 2.7)
    if img:
        story.append(Spacer(1, 0.15 * inch))
        story.append(img)
        story.append(Paragraph('Gráfico 4 — Cumplimiento temporal: formularios Dentro del PAC evaluados contra su fecha planificada.', est['Leyenda']))
    img = _pdf_imagen(grafico_barras_ejecucion_mensual_planer(d['temporal_mensual_planer']['meses']), 5.6, 2.4)
    if img:
        story.append(img)
        story.append(Paragraph(f'Gráfico 5 — Ejecución mensual del Plan de Compras {anho} (fichas Ejecutadas/Pendientes/Atrasadas).', est['Leyenda']))
    story.append(PageBreak())

    # --- Metodología -----------------------------------------------------------
    story.append(Paragraph('Metodología', est['TituloCapitulo']))
    for parrafo in _PARRAFOS_METODOLOGIA:
        story.append(Paragraph(parrafo, est['Cuerpo']))
    story.append(PageBreak())

    # --- Capítulos por Subdirección ----------------------------------------------
    subs_institucionales = [s for s in d['subdirecciones'] if s['institucional']]
    otras_ramas = [s for s in d['subdirecciones'] if not s['institucional']]

    for sub in subs_institucionales:
        nombre_display = _nombre_subdireccion_display(sub['nombre'])
        story.append(Paragraph(nombre_display, est['TituloCapitulo']))

        story.append(Paragraph(f'Formularios de Solicitud de Compra — Período {label}', est['TituloSeccion']))
        if sub['fsc'] and sub['fsc']['total']:
            mejor_rk, peor_rk = _mejor_peor_de_ranking(d['rankings_depto'], sub['nombre'])
            story.append(Paragraph(parrafo_capitulo_subdireccion(nombre_display, sub['fsc'], mejor_rk, peor_rk), est['Cuerpo']))
            if sub['fsc']['departamentos']:
                story.append(_pdf_tabla_departamentos(sub['fsc']['departamentos']))

            img = _pdf_imagen(grafico_donut_dentro_fuera(sub['fsc']['pct_dentro'], titulo=f'{nombre_display} — Dentro/Fuera PAC'), 2.4, 2.4)
            if img:
                story.append(Spacer(1, 0.1 * inch))
                story.append(img)
                story.append(Paragraph(
                    f'Distribución de los {_n(sub["fsc"]["total"])} formularios de {nombre_display} entre Dentro y '
                    'Fuera del PAC durante el período.', est['Leyenda'],
                ))

            if sub['formularios_detalle']:
                story.append(Paragraph(
                    f'Detalle de los {_n(len(sub["formularios_detalle"]))} formularios individuales de {nombre_display} '
                    'en el período:', est['Cuerpo'],
                ))
                story.extend(_pdf_tabla_formularios_detalle(sub['formularios_detalle']))
        else:
            story.append(Paragraph('Sin formularios registrados en el período para esta subdirección.', est['Cuerpo']))

        story.append(Spacer(1, 0.12 * inch))
        story.append(Paragraph(f'Ejecución del Plan de Compras — Año PAC {anho}', est['TituloSeccion']))
        if sub['planer'] and sub['planer']['total']:
            p = sub['planer']
            story.append(Paragraph(
                f"{_n(p['total'])} fichas del Plan de Compras — {p['pct_ejecutado']}% ejecutadas "
                f"({_n(p['ejecutados'])} ejecutadas, {_n(p['pendientes'])} pendientes, {_n(p['atrasados'])} atrasadas). "
                f"Monto total planificado: {_money(p['monto_total'])}.",
                est['Cuerpo'],
            ))
            if p['departamentos']:
                story.append(_pdf_tabla_deptos_ejecucion(p['departamentos']))
            if sub['responsables']:
                story.append(Paragraph('Responsables: ' + ', '.join(sub['responsables']), est['Metadato']))

            img = _pdf_imagen(
                grafico_donut_ejecucion_planer(p['ejecutados'], p['pendientes'], p['atrasados'], titulo=f'{nombre_display} — Ejecución del Plan de Compras'),
                2.7, 2.5,
            )
            if img:
                story.append(Spacer(1, 0.1 * inch))
                story.append(img)
                story.append(Paragraph(
                    f'Estado de ejecución de los {_n(p["total"])} proyectos del Plan de Compras {anho} de {nombre_display}.',
                    est['Leyenda'],
                ))

            if sub['fichas_detalle']:
                story.append(Paragraph(
                    f'Detalle de los {_n(len(sub["fichas_detalle"]))} proyectos del Plan de Compras {anho} de '
                    f'{nombre_display}, indicando con qué formulario(s) se ejecutó cada uno:', est['Cuerpo'],
                ))
                story.extend(_pdf_tabla_fichas_ejecucion_detalle(sub['fichas_detalle']))
        else:
            story.append(Paragraph('Sin fichas del Plan de Compras registradas para esta subdirección.', est['Cuerpo']))
        story.append(PageBreak())

    # --- Rankings Institucionales -------------------------------------------
    story.append(Paragraph('Rankings Institucionales', est['TituloCapitulo']))
    story.append(Paragraph(
        'El score compuesto pondera 40% el % Dentro del PAC, 40% el % de cumplimiento temporal y 20% el % '
        'Dentro del PAC ponderado por monto.', est['Cuerpo'],
    ))
    story.append(Paragraph('Mejores departamentos', est['TituloSeccion']))
    if d['rankings_depto']['mejores']:
        story.append(_pdf_tabla_ranking_depto(d['rankings_depto']['mejores']))
    else:
        story.append(Paragraph('Sin datos suficientes para generar el ranking en este período.', est['Cuerpo']))
    story.append(Spacer(1, 0.15 * inch))
    story.append(Paragraph('Departamentos con mayor oportunidad de mejora', est['TituloSeccion']))
    if d['rankings_depto']['peores']:
        story.append(_pdf_tabla_ranking_depto(d['rankings_depto']['peores']))
    else:
        story.append(Paragraph('Sin datos suficientes para generar el ranking en este período.', est['Cuerpo']))
    story.append(PageBreak())

    story.append(Paragraph('Mejores formularios', est['TituloSeccion']))
    if d['rankings_formulario']['mejores']:
        story.append(_pdf_tabla_ranking_formulario(d['rankings_formulario']['mejores']))
    else:
        story.append(Paragraph('Sin datos suficientes para generar el ranking en este período.', est['Cuerpo']))
    story.append(Spacer(1, 0.15 * inch))
    story.append(Paragraph('Formularios con mayor oportunidad de mejora', est['TituloSeccion']))
    if d['rankings_formulario']['peores']:
        story.append(_pdf_tabla_ranking_formulario(d['rankings_formulario']['peores']))
    else:
        story.append(Paragraph('Sin datos suficientes para generar el ranking en este período.', est['Cuerpo']))
    story.append(PageBreak())

    # --- Alertas y Seguimiento ------------------------------------------------
    story.append(Paragraph('Alertas y Seguimiento', est['TituloCapitulo']))
    story.append(Paragraph(
        'Las secciones siguientes se agrupan por subdirección institucional para facilitar el seguimiento y '
        'la coordinación con cada responsable.', est['Cuerpo'],
    ))
    story.append(Paragraph('Proyectos planificados sin ningún formulario Dentro PAC derivado todavía', est['TituloSeccion']))
    _render_grupos_subdireccion_pdf(
        story, d['proyectos_sin_iniciar'], _pdf_tabla_proyectos_sin_iniciar, est,
        'No hay proyectos planificados sin iniciar en el período.',
    )

    story.append(Paragraph(f'Fichas a ejecutar el próximo mes (Año PAC {anho})', est['TituloSeccion']))
    _render_grupos_subdireccion_pdf(
        story, d['proximo_mes'], _pdf_tabla_fichas_pac, est,
        'No hay fichas planificadas para el próximo mes.',
    )

    story.append(Paragraph(f'Fichas atrasadas (Año PAC {anho})', est['TituloSeccion']))
    _render_grupos_subdireccion_pdf(
        story, d['atrasadas'], _pdf_tabla_fichas_pac, est,
        'No hay fichas atrasadas en el período.',
    )
    story.append(PageBreak())

    # --- Conclusiones y Recomendaciones ---------------------------------------
    story.append(Paragraph('Conclusiones y Recomendaciones', est['TituloCapitulo']))
    story.append(Paragraph(parrafo_conclusiones(label, d['dentro_fuera']['kpis'], d['temporal']['kpis']), est['Cuerpo']))
    story.append(Paragraph(
        f'En cuanto a la ejecución del Plan de Compras {anho}, {d["pct_ejecutado_fichas"]}% de las fichas '
        'planificadas cuentan hoy con formulario u orden de compra enlazada. '
        + ('Se recomienda priorizar la revisión de las fichas atrasadas listadas en Alertas y Seguimiento y '
           'coordinar con los responsables de cada subdirección la regularización de los proyectos pendientes.'
           if d['pct_ejecutado_fichas'] < 70 else
           'El nivel de ejecución es satisfactorio; se recomienda mantener el seguimiento mensual para sostener '
           'la tendencia.'),
        est['Cuerpo'],
    ))
    story.append(PageBreak())

    # --- Anexo: otras unidades de la red ----------------------------------------
    story.append(Paragraph('Anexo — Otras Unidades de la Red', est['TituloCapitulo']))
    story.append(Paragraph(
        'Detalle abreviado de las restantes unidades de la red asistencial (fuera de la Dirección Servicio de '
        'Salud Osorno) y de los formularios que no pudieron clasificarse organizacionalmente.', est['Cuerpo'],
    ))
    if otras_ramas:
        for sub in otras_ramas:
            story.append(Paragraph(_nombre_subdireccion_display(sub['nombre']), est['TituloSeccion']))
            if sub['fsc'] and sub['fsc']['departamentos']:
                story.append(_pdf_tabla_departamentos(sub['fsc']['departamentos']))
            if sub['planer'] and sub['planer']['departamentos']:
                story.append(Spacer(1, 0.08 * inch))
                story.append(Paragraph('Ejecución del Plan de Compras:', est['Cuerpo']))
                story.append(_pdf_tabla_deptos_ejecucion(sub['planer']['departamentos']))
            story.append(Spacer(1, 0.1 * inch))
    else:
        story.append(Paragraph('Sin registros de otras unidades de la red en este período.', est['Cuerpo']))

    doc.multiBuild(story)
    buf.seek(0)
    return buf


# =============================================================================
# PPT — presentación institucional completa (decisión del usuario 2026-07-21,
# ampliada el mismo día): parte con una sección general (KPIs, comparativas,
# evolución mensual) y luego desarrolla UN capítulo por subdirección con sus
# propios KPIs, gráficos y tablas por departamento — pensada para poder
# exponerse completa o subdirección por subdirección en una reunión con Alta
# Dirección. Ya NO es un resumen de una sola diapositiva por subdirección
# (versión anterior, más liviana); ese nivel de detalle ahora vive en
# `_ppt_capitulo_subdireccion`, espejo del capítulo de Word/PDF pero en slides.
# =============================================================================

# Paleta y medidas — extraídas de `PAC_SSO_2026 (1).pptx` (mazo institucional de
# referencia del Subdepartamento de Abastecimiento), para que los informes
# generados por el sistema tengan el mismo lenguaje visual. Ver detalle de
# patrones de diapositiva en el plan de implementación.
COLOR_TITULO_PPT = PptxRGBColor(0x00, 0x23, 0x4e)          # navy primario
COLOR_ACENTO_PPT = PptxRGBColor(0x4d, 0xaa, 0xdf)           # azul acento (barras finas)
COLOR_TEXTO_PPT = PptxRGBColor(0x33, 0x33, 0x33)            # texto cuerpo
COLOR_PIE_PPT = PptxRGBColor(0xaa, 0xaa, 0xaa)              # texto de pie (slides claras)

COLOR_PPT_AZUL_PROFUNDO = PptxRGBColor(0x00, 0x4b, 0x8d)
COLOR_PPT_AZUL_MEDIO = PptxRGBColor(0x00, 0x5a, 0x9e)
COLOR_PPT_FONDO_CLARO = PptxRGBColor(0xf0, 0xf5, 0xfa)
COLOR_PPT_FILA_ALT = PptxRGBColor(0xf8, 0xfa, 0xff)
COLOR_PPT_TEXTO_MUTED = PptxRGBColor(0x64, 0x74, 0x8b)
COLOR_PPT_TEXTO_FORMULA = PptxRGBColor(0x1e, 0x3a, 0x5f)
COLOR_PPT_VERDE = PptxRGBColor(0x16, 0xa3, 0x4a)
COLOR_PPT_AMBAR = PptxRGBColor(0xd9, 0x77, 0x06)
COLOR_PPT_ROJO = PptxRGBColor(0xdc, 0x26, 0x26)
COLOR_PPT_DESTACADO_BG = PptxRGBColor(0xeb, 0xf2, 0xff)
COLOR_PPT_NUMERO_DIVISOR = PptxRGBColor(0x1a, 0x3a, 0x6b)   # número gigante de separador de sección
COLOR_PPT_PIE_OSCURO = PptxRGBColor(0xb0, 0xc8, 0xe0)       # texto de pie sobre fondo navy

# Lienzo 16:9 estándar de PowerPoint (13.333x7.5") — el alto no cambia respecto
# al 4:3 anterior (10x7.5"), así que la mayoría de los `top=`/`height=` ya
# existentes en este archivo siguen siendo válidos; solo se reescalaron los
# `left=`/`width=` para aprovechar el ancho nuevo.
PPT_ANCHO = 13.333
PPT_ALTO = 7.5
PPT_MARGEN = 0.5
PPT_ANCHO_CONTENIDO = PPT_ANCHO - 2 * PPT_MARGEN  # ≈ 12.33"


def _ppt_slide_en_blanco(prs):
    """Slide de CONTENIDO (fondo claro): fondo `F0F5FA` + barra superior navy
    delgada, mimetizando las diapositivas de contenido del mazo de referencia.
    El pie (barra inferior + texto + numeración) se agrega al final vía
    `_ppt_pie_pagina`/`_ppt_numerar_diapositivas` — no hace falta tocar cada
    call-site que ya usa esta función."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_PPT_FONDO_CLARO
    barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0), PptxInches(PPT_ANCHO), PptxInches(0.08))
    barra.fill.solid()
    barra.fill.fore_color.rgb = COLOR_TITULO_PPT
    barra.line.fill.background()
    barra.shadow.inherit = False
    return slide


def _ppt_slide_oscura(prs):
    """Slide OSCURA (fondo navy) para portadas y separadores de sección: barras
    de acento superior/inferior + barra vertical izquierda, como en el mazo de
    referencia. El contenido (título, número, etc.) se agrega por el llamador."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_TITULO_PPT
    for top in (0, PPT_ALTO - 0.08):
        barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(top), PptxInches(PPT_ANCHO), PptxInches(0.08))
        barra.fill.solid()
        barra.fill.fore_color.rgb = COLOR_ACENTO_PPT
        barra.line.fill.background()
        barra.shadow.inherit = False
    lateral = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0.08), PptxInches(0.3), PptxInches(PPT_ALTO - 0.16))
    lateral.fill.solid()
    lateral.fill.fore_color.rgb = COLOR_PPT_AZUL_PROFUNDO
    lateral.line.fill.background()
    lateral.shadow.inherit = False
    return slide


def _ppt_titulo(slide, texto, subtitulo=None, top=PptxInches(0.2), tamano=24):
    caja = slide.shapes.add_textbox(PptxInches(PPT_MARGEN), top, PptxInches(PPT_ANCHO_CONTENIDO), PptxInches(0.55))
    p = caja.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = texto
    run.font.size = PptxPt(tamano)
    run.font.bold = True
    run.font.color.rgb = COLOR_TITULO_PPT
    if subtitulo:
        caja_sub = slide.shapes.add_textbox(
            PptxInches(PPT_MARGEN), top + PptxInches(0.5), PptxInches(PPT_ANCHO_CONTENIDO), PptxInches(0.3),
        )
        p_sub = caja_sub.text_frame.paragraphs[0]
        run_sub = p_sub.add_run()
        run_sub.text = subtitulo
        run_sub.font.size = PptxPt(11)
        run_sub.font.color.rgb = COLOR_PPT_TEXTO_MUTED
    return caja


def _ppt_parrafo(slide, texto, top, left=PptxInches(PPT_MARGEN), width=PptxInches(PPT_ANCHO_CONTENIDO), tamano=14, height=PptxInches(1.2)):
    caja = slide.shapes.add_textbox(left, top, width, height)
    tf = caja.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = texto
    run.font.size = PptxPt(tamano)
    run.font.color.rgb = COLOR_TEXTO_PPT
    return caja


def _ppt_imagen(slide, img_bytes, left, top, width):
    if img_bytes:
        slide.shapes.add_picture(img_bytes, left, top, width=width)


def _ppt_es_slide_oscura(slide):
    try:
        return slide.background.fill.fore_color.rgb == COLOR_TITULO_PPT
    except Exception:
        return False


def _ppt_pie_pagina(slide):
    """Barra inferior navy + texto de pie — solo en slides de contenido (fondo
    claro). Las slides oscuras (portada/separadores) ya dibujan su propio pie
    como parte de su contenido custom, así que se ignoran aquí — permite seguir
    llamando esta función sobre TODAS las slides sin distinción."""
    if _ppt_es_slide_oscura(slide):
        return
    barra = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(PPT_ALTO - 0.08), PptxInches(PPT_ANCHO), PptxInches(0.08),
    )
    barra.fill.solid()
    barra.fill.fore_color.rgb = COLOR_TITULO_PPT
    barra.line.fill.background()
    barra.shadow.inherit = False
    caja = slide.shapes.add_textbox(PptxInches(PPT_MARGEN), PptxInches(PPT_ALTO - 0.34), PptxInches(PPT_ANCHO_CONTENIDO - 0.7), PptxInches(0.24))
    p = caja.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = NOMBRE_INSTITUCION
    run.font.size = PptxPt(8)
    run.font.color.rgb = COLOR_PIE_PPT


def _ppt_numerar_diapositivas(prs):
    total = len(prs.slides._sldIdLst)
    for i, slide in enumerate(prs.slides, start=1):
        if _ppt_es_slide_oscura(slide):
            continue
        caja = slide.shapes.add_textbox(
            PptxInches(PPT_ANCHO - PPT_MARGEN - 0.6), PptxInches(PPT_ALTO - 0.34), PptxInches(0.6), PptxInches(0.24),
        )
        p = caja.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f'{i}/{total}'
        run.font.size = PptxPt(8)
        run.font.color.rgb = COLOR_PIE_PPT


def _ppt_portada(prs, badge_texto, titulo_lineas, subtitulo_inst, meta_texto, logo_bytes=None, imagen_bytes=None):
    """Portada estilo `PAC_SSO_2026.pptx`: fondo navy, badge azul acento, título
    grande blanco (multilínea), subtítulo institucional, línea divisoria fina,
    meta-texto de pie. `titulo_lineas` es una lista de strings (una por línea)."""
    slide = _ppt_slide_oscura(prs)
    if logo_bytes:
        slide.shapes.add_picture(logo_bytes, PptxInches(PPT_ANCHO - PPT_MARGEN - 1.3), PptxInches(0.35), height=PptxInches(1.0))
    if badge_texto:
        badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0.7), PptxInches(0.5), PptxInches(3.4), PptxInches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_ACENTO_PPT
        badge.line.fill.background()
        badge.shadow.inherit = False
        tf = badge.text_frame
        tf.margin_left = PptxInches(0.08)
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = badge_texto
        run.font.size = PptxPt(13)
        run.font.bold = True
        run.font.color.rgb = COLOR_TITULO_PPT
    caja_t = slide.shapes.add_textbox(PptxInches(0.7), PptxInches(1.35), PptxInches(PPT_ANCHO - 1.4), PptxInches(2.2))
    tf_t = caja_t.text_frame
    tf_t.word_wrap = True
    for i, linea in enumerate(titulo_lineas):
        p = tf_t.paragraphs[0] if i == 0 else tf_t.add_paragraph()
        run = p.add_run()
        run.text = linea
        run.font.size = PptxPt(34)
        run.font.bold = True
        run.font.color.rgb = PptxRGBColor(0xff, 0xff, 0xff)
    if subtitulo_inst:
        caja_sub = slide.shapes.add_textbox(PptxInches(0.7), PptxInches(1.35 + 0.7 * len(titulo_lineas) + 0.3), PptxInches(PPT_ANCHO - 1.4), PptxInches(0.6))
        p_sub = caja_sub.text_frame.paragraphs[0]
        run_sub = p_sub.add_run()
        run_sub.text = subtitulo_inst
        run_sub.font.size = PptxPt(20)
        run_sub.font.color.rgb = COLOR_ACENTO_PPT
    if imagen_bytes:
        # Ancho acotado (no solo alto) para que quepa entre el subtítulo y la línea
        # divisoria del pie sin desbordar el lienzo — la imagen del edificio
        # institucional tiene una relación de aspecto más alta que ancha.
        slide.shapes.add_picture(imagen_bytes, PptxInches((PPT_ANCHO - 3.5) / 2), PptxInches(3.75), width=PptxInches(3.5))
    if meta_texto:
        divisor = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0.7), PptxInches(PPT_ALTO - 1.0), PptxInches(6), PptxInches(0.035))
        divisor.fill.solid()
        divisor.fill.fore_color.rgb = COLOR_ACENTO_PPT
        divisor.line.fill.background()
        divisor.shadow.inherit = False
        caja_meta = slide.shapes.add_textbox(PptxInches(0.7), PptxInches(PPT_ALTO - 0.85), PptxInches(PPT_ANCHO - 1.4), PptxInches(0.6))
        p_meta = caja_meta.text_frame.paragraphs[0]
        run_meta = p_meta.add_run()
        run_meta.text = meta_texto
        run_meta.font.size = PptxPt(13)
        run_meta.font.color.rgb = COLOR_PPT_PIE_OSCURO
    return slide


def _ppt_divisor_seccion(prs, numero, titulo_lineas, subtitulo=None):
    """Separador de sección estilo `PAC_SSO_2026.pptx`: fondo navy, número
    gigante, título multilínea blanco bold, subtítulo azul acento."""
    slide = _ppt_slide_oscura(prs)
    caja_n = slide.shapes.add_textbox(PptxInches(0.7), PptxInches(0.6), PptxInches(3), PptxInches(1.5))
    p_n = caja_n.text_frame.paragraphs[0]
    run_n = p_n.add_run()
    run_n.text = numero
    run_n.font.size = PptxPt(80)
    run_n.font.bold = True
    run_n.font.color.rgb = COLOR_PPT_NUMERO_DIVISOR
    caja_t = slide.shapes.add_textbox(PptxInches(0.7), PptxInches(2.3), PptxInches(PPT_ANCHO - 1.4), PptxInches(2.2))
    tf_t = caja_t.text_frame
    tf_t.word_wrap = True
    for i, linea in enumerate(titulo_lineas):
        p = tf_t.paragraphs[0] if i == 0 else tf_t.add_paragraph()
        run = p.add_run()
        run.text = linea
        run.font.size = PptxPt(36)
        run.font.bold = True
        run.font.color.rgb = PptxRGBColor(0xff, 0xff, 0xff)
    if subtitulo:
        caja_sub = slide.shapes.add_textbox(
            PptxInches(0.7), PptxInches(2.3 + 0.75 * len(titulo_lineas) + 0.2), PptxInches(PPT_ANCHO - 1.4), PptxInches(0.6),
        )
        p_sub = caja_sub.text_frame.paragraphs[0]
        run_sub = p_sub.add_run()
        run_sub.text = subtitulo
        run_sub.font.size = PptxPt(15)
        run_sub.font.color.rgb = COLOR_ACENTO_PPT
    return slide


def _ppt_tarjeta_kpi(slide, left, top, width, height, valor, etiqueta, color_bg=None, color_texto_etiqueta=None):
    """Tile KPI sólido (número grande blanco + etiqueta) — como los stat cards
    de las diapositivas 15/16 de la referencia."""
    color_bg = color_bg or COLOR_TITULO_PPT
    color_texto_etiqueta = color_texto_etiqueta or PptxRGBColor(0xcc, 0xe0, 0xf5)
    tile = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    tile.fill.solid()
    tile.fill.fore_color.rgb = color_bg
    tile.line.fill.background()
    tile.shadow.inherit = False
    caja_v = slide.shapes.add_textbox(left + PptxInches(0.06), top + PptxInches(0.04), width - PptxInches(0.12), height * 0.62)
    p_v = caja_v.text_frame.paragraphs[0]
    run_v = p_v.add_run()
    run_v.text = str(valor)
    run_v.font.size = PptxPt(24)
    run_v.font.bold = True
    run_v.font.color.rgb = PptxRGBColor(0xff, 0xff, 0xff)
    caja_e = slide.shapes.add_textbox(left + PptxInches(0.06), top + height * 0.62, width - PptxInches(0.12), height * 0.36)
    p_e = caja_e.text_frame.paragraphs[0]
    run_e = p_e.add_run()
    run_e.text = etiqueta
    run_e.font.size = PptxPt(10.5)
    run_e.font.color.rgb = color_texto_etiqueta


def _ppt_fila_indicador(slide, top, etiqueta, valor, color, comparacion=None, explicacion=None, flecha='→'):
    """Fila de indicador estilo diapositiva 13 de la referencia: icono de
    tendencia + etiqueta + valor grande coloreado + comparación gris +
    explicación en texto oscuro."""
    alto = PptxInches(0.66)
    fondo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(PPT_MARGEN), top, PptxInches(PPT_ANCHO_CONTENIDO), alto)
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = PptxRGBColor(0xff, 0xff, 0xff)
    fondo.line.fill.background()
    fondo.shadow.inherit = False
    barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(PPT_MARGEN), top, PptxInches(0.09), alto)
    barra.fill.solid()
    barra.fill.fore_color.rgb = color
    barra.line.fill.background()
    barra.shadow.inherit = False
    caja_icono = slide.shapes.add_textbox(PptxInches(PPT_MARGEN + 0.18), top + PptxInches(0.13), PptxInches(0.4), PptxInches(0.4))
    p_icono = caja_icono.text_frame.paragraphs[0]
    run_icono = p_icono.add_run()
    run_icono.text = flecha
    run_icono.font.size = PptxPt(18)
    run_icono.font.bold = True
    run_icono.font.color.rgb = color
    caja_e = slide.shapes.add_textbox(PptxInches(PPT_MARGEN + 0.75), top + PptxInches(0.06), PptxInches(2.7), PptxInches(0.3))
    p_e = caja_e.text_frame.paragraphs[0]
    run_e = p_e.add_run()
    run_e.text = etiqueta
    run_e.font.size = PptxPt(12)
    run_e.font.bold = True
    run_e.font.color.rgb = COLOR_TITULO_PPT
    caja_v = slide.shapes.add_textbox(PptxInches(PPT_MARGEN + 0.75), top + PptxInches(0.32), PptxInches(1.6), PptxInches(0.3))
    p_v = caja_v.text_frame.paragraphs[0]
    run_v = p_v.add_run()
    run_v.text = str(valor)
    run_v.font.size = PptxPt(16)
    run_v.font.bold = True
    run_v.font.color.rgb = color
    if comparacion:
        caja_c = slide.shapes.add_textbox(PptxInches(PPT_MARGEN + 2.5), top + PptxInches(0.35), PptxInches(2.8), PptxInches(0.25))
        p_c = caja_c.text_frame.paragraphs[0]
        run_c = p_c.add_run()
        run_c.text = comparacion
        run_c.font.size = PptxPt(9.5)
        run_c.font.color.rgb = COLOR_PPT_TEXTO_MUTED
    if explicacion:
        caja_x = slide.shapes.add_textbox(PptxInches(PPT_MARGEN + 5.5), top + PptxInches(0.08), PptxInches(PPT_ANCHO_CONTENIDO - 5.6), alto - PptxInches(0.16))
        tf_x = caja_x.text_frame
        tf_x.word_wrap = True
        p_x = tf_x.paragraphs[0]
        run_x = p_x.add_run()
        run_x.text = explicacion
        run_x.font.size = PptxPt(9.5)
        run_x.font.color.rgb = COLOR_TEXTO_PPT


def _ppt_tarjeta_numerada(slide, top, numero, titulo, descripcion=None, color=None, height=0.75):
    """Tarjeta blanca con badge numerado — Agenda / listas de conclusiones,
    como las diapositivas 2 y 8 de la referencia."""
    color = color or COLOR_TITULO_PPT
    alto = PptxInches(height)
    fondo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(PPT_MARGEN), top, PptxInches(PPT_ANCHO_CONTENIDO), alto)
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = PptxRGBColor(0xff, 0xff, 0xff)
    fondo.line.fill.background()
    fondo.shadow.inherit = False
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(PPT_MARGEN), top, PptxInches(0.62), alto)
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    badge.shadow.inherit = False
    tf_b = badge.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    run_b = p_b.add_run()
    run_b.text = numero
    run_b.font.size = PptxPt(18)
    run_b.font.bold = True
    run_b.font.color.rgb = PptxRGBColor(0xff, 0xff, 0xff)
    caja_t = slide.shapes.add_textbox(PptxInches(PPT_MARGEN + 0.78), top + PptxInches(0.06), PptxInches(PPT_ANCHO_CONTENIDO - 0.9), PptxInches(0.32))
    p_t = caja_t.text_frame.paragraphs[0]
    run_t = p_t.add_run()
    run_t.text = titulo
    run_t.font.size = PptxPt(14)
    run_t.font.bold = True
    run_t.font.color.rgb = COLOR_TITULO_PPT
    if descripcion:
        caja_d = slide.shapes.add_textbox(PptxInches(PPT_MARGEN + 0.78), top + PptxInches(0.38), PptxInches(PPT_ANCHO_CONTENIDO - 0.9), alto - PptxInches(0.42))
        tf_d = caja_d.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        run_d = p_d.add_run()
        run_d.text = descripcion
        run_d.font.size = PptxPt(10.5)
        run_d.font.color.rgb = COLOR_PPT_TEXTO_MUTED


def _ppt_estilizar_tabla(tabla, color_header=None, zebra=True, resaltar_ultima_fila=False, tamano_fuente=11):
    """Aplica el estilo de tabla del mazo de referencia: header navy con texto
    blanco bold, filas zebra (blanco/`F8FAFF`), fila TOTAL opcional resaltada.
    Se llama al final de cada función `_ppt_tabla_*`, sobre una tabla ya
    poblada con texto."""
    color_header = color_header or COLOR_TITULO_PPT
    n_filas = len(tabla.rows)
    for r, row in enumerate(tabla.rows):
        es_header = r == 0
        es_total = resaltar_ultima_fila and r == n_filas - 1
        for cell in row.cells:
            cell.margin_top = PptxInches(0.03)
            cell.margin_bottom = PptxInches(0.03)
            cell.fill.solid()
            if es_header:
                cell.fill.fore_color.rgb = color_header
            elif es_total:
                cell.fill.fore_color.rgb = COLOR_PPT_DESTACADO_BG
            elif zebra and r % 2 == 0:
                cell.fill.fore_color.rgb = COLOR_PPT_FILA_ALT
            else:
                cell.fill.fore_color.rgb = PptxRGBColor(0xff, 0xff, 0xff)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = PptxPt(tamano_fuente)
                    if es_header:
                        run.font.bold = True
                        run.font.color.rgb = PptxRGBColor(0xff, 0xff, 0xff)
                    elif es_total:
                        run.font.bold = True
                        run.font.color.rgb = COLOR_TITULO_PPT
                    else:
                        run.font.color.rgb = COLOR_TEXTO_PPT


def _ppt_tabla_resumen_combinado(slide, subdirecciones, top):
    """Una sola tabla que resume las 4 subdirecciones institucionales — versión
    ejecutiva del capítulo por subdirección exhaustivo de Word/PDF."""
    filas_datos = [s for s in subdirecciones if s['institucional']]
    n_filas = len(filas_datos) + 1
    tabla_shape = slide.shapes.add_table(n_filas, 5, PptxInches(PPT_MARGEN), top, PptxInches(PPT_ANCHO_CONTENIDO), PptxInches(0.45 * n_filas))
    tabla = tabla_shape.table
    for i, h in enumerate(['Subdirección', 'Formularios', '% Dentro PAC', 'Fichas Plan Compras', '% Ejecutado']):
        tabla.cell(0, i).text = h
    for r, s in enumerate(filas_datos, start=1):
        fsc, planer = s['fsc'], s['planer']
        tabla.cell(r, 0).text = _nombre_subdireccion_display(s['nombre'])
        tabla.cell(r, 1).text = str(fsc['total']) if fsc else '—'
        tabla.cell(r, 2).text = f"{fsc['pct_dentro']}%" if fsc else '—'
        tabla.cell(r, 3).text = str(planer['total']) if planer else '—'
        tabla.cell(r, 4).text = f"{planer['pct_ejecutado']}%" if planer else '—'
    _ppt_estilizar_tabla(tabla, tamano_fuente=12)


def _ppt_tabla_ranking(slide, filas, top, limite=5):
    filas_datos = filas[:limite]
    n_filas = len(filas_datos) + 1
    tabla_shape = slide.shapes.add_table(n_filas, 4, PptxInches(PPT_MARGEN), top, PptxInches(PPT_ANCHO_CONTENIDO), PptxInches(0.42 * n_filas))
    tabla = tabla_shape.table
    for i, h in enumerate(['Departamento', 'Subdirección', 'Score', '% Dentro']):
        tabla.cell(0, i).text = h
    for r, f in enumerate(filas_datos, start=1):
        tabla.cell(r, 0).text = f['nombre'].title()
        tabla.cell(r, 1).text = _nombre_subdireccion_display(f['subdireccion'])
        tabla.cell(r, 2).text = f"{f['score']:.0f}"
        tabla.cell(r, 3).text = f"{f['pct_dentro']}%"
    _ppt_estilizar_tabla(tabla, tamano_fuente=12)
    if len(filas) > limite:
        _ppt_parrafo(
            slide, f'… y {len(filas) - limite} departamento(s) adicional(es) — ver el ranking completo en el informe Word/PDF.',
            top=top + PptxInches(0.42 * n_filas) + PptxInches(0.1), tamano=10, height=PptxInches(0.4),
        )


def _ppt_tabla_ranking_formulario(slide, filas, top, limite=5):
    """Ranking de formularios individuales — versión slide de `_tabla_ranking_formulario`
    (Word). Ausente del PPT hasta la revisión de código 2026-07-27: el PPT solo mostraba
    el ranking por departamento, omitiendo por completo la dimensión por formulario que
    Word/PDF sí incluyen."""
    filas_datos = filas[:limite]
    n_filas = len(filas_datos) + 1
    tabla_shape = slide.shapes.add_table(n_filas, 5, PptxInches(PPT_MARGEN), top, PptxInches(PPT_ANCHO_CONTENIDO), PptxInches(0.42 * n_filas))
    tabla = tabla_shape.table
    _ppt_fijar_anchos_columnas(tabla, [1.47, 4.55, 1.47, 2.27, 2.54])
    for i, h in enumerate(['Folio/Año', 'Unidad Requirente', 'Dentro/Fuera', 'Monto', 'Score']):
        tabla.cell(0, i).text = h
    for r, f in enumerate(filas_datos, start=1):
        tabla.cell(r, 0).text = f"{f['folio']}/{f['anho']}"
        tabla.cell(r, 1).text = _truncar(f['unidad_requirente'], 45)
        tabla.cell(r, 2).text = 'Dentro' if f['dentro_fuera_pac'] == 'DENTRO' else 'Fuera'
        tabla.cell(r, 3).text = _money(f['monto_estimado'])
        tabla.cell(r, 4).text = f"{f['score']:.0f}"
    _ppt_estilizar_tabla(tabla, tamano_fuente=12)
    if len(filas) > limite:
        _ppt_parrafo(
            slide, f'… y {len(filas) - limite} formulario(s) adicional(es) — ver el ranking completo en el informe Word/PDF.',
            top=top + PptxInches(0.42 * n_filas) + PptxInches(0.1), tamano=10, height=PptxInches(0.4),
        )


def _ppt_fijar_anchos_columnas(tabla, anchos_pulgadas):
    """Fija el ancho de TODAS las columnas explícitamente. python-pptx no reajusta
    las demás columnas cuando se cambia el ancho de solo una — el ancho total de la
    tabla pasa a ser la suma de columnas, pudiendo salirse del borde de la diapositiva
    si solo se agranda una y las otras quedan en su ancho parejo original (bug real
    detectado 2026-07-23: una tabla de 9.2" declaradas terminaba en 10.76" — 0.76"
    fuera de una diapositiva de 10"). Siempre fijar todas para que la suma coincida
    con el ancho total declarado en `add_table`."""
    for col, ancho in zip(tabla.columns, anchos_pulgadas):
        col.width = PptxInches(ancho)


def _ppt_tabla_depto_dentro_fuera(slide, departamentos, top, limite=12):
    """Tabla Dentro/Fuera PAC por departamento — capítulo por subdirección del PPT.
    Mismos campos que `_tabla_departamentos` (Word), formato compacto de slide."""
    filas_datos = departamentos[:limite]
    n_filas = len(filas_datos) + 1
    tabla_shape = slide.shapes.add_table(n_filas, 5, PptxInches(PPT_MARGEN), top, PptxInches(PPT_ANCHO_CONTENIDO), PptxInches(0.4 * n_filas))
    tabla = tabla_shape.table
    _ppt_fijar_anchos_columnas(tabla, [4.55, 1.94, 1.94, 1.94, 1.94])
    for i, h in enumerate(['Departamento', 'Total', 'Dentro', 'Fuera', '% Dentro']):
        tabla.cell(0, i).text = h
    for r, d in enumerate(filas_datos, start=1):
        tabla.cell(r, 0).text = _truncar(d['nombre'].title(), 42)
        tabla.cell(r, 1).text = str(d['total'])
        tabla.cell(r, 2).text = str(d['dentro'])
        tabla.cell(r, 3).text = str(d['fuera'])
        tabla.cell(r, 4).text = f"{d['pct_dentro']}%"
    _ppt_estilizar_tabla(tabla, tamano_fuente=12)
    if len(departamentos) > limite:
        _ppt_parrafo(
            slide, f'… y {len(departamentos) - limite} departamento(s) adicional(es) — ver detalle en el informe Word/PDF.',
            top=top + PptxInches(0.4 * n_filas) + PptxInches(0.1), tamano=10, height=PptxInches(0.4),
        )


def _ppt_tabla_depto_ejecucion(slide, departamentos, top, limite=12):
    """Tabla Ejecución del Plan de Compras por departamento — capítulo por
    subdirección del PPT. Mismos campos que `_tabla_deptos_ejecucion` (Word)."""
    filas_datos = departamentos[:limite]
    n_filas = len(filas_datos) + 1
    tabla_shape = slide.shapes.add_table(n_filas, 5, PptxInches(PPT_MARGEN), top, PptxInches(PPT_ANCHO_CONTENIDO), PptxInches(0.4 * n_filas))
    tabla = tabla_shape.table
    _ppt_fijar_anchos_columnas(tabla, [4.55, 1.94, 1.94, 1.94, 1.94])
    for i, h in enumerate(['Departamento', 'Total', 'Ejecutadas', 'Atrasadas', '% Ejecutado']):
        tabla.cell(0, i).text = h
    for r, d in enumerate(filas_datos, start=1):
        tabla.cell(r, 0).text = _truncar(d['nombre'].title(), 42)
        tabla.cell(r, 1).text = str(d['total'])
        tabla.cell(r, 2).text = str(d['ejecutados'])
        tabla.cell(r, 3).text = str(d['atrasados'])
        tabla.cell(r, 4).text = f"{d['pct_ejecutado']}%"
    _ppt_estilizar_tabla(tabla, tamano_fuente=12)
    if len(departamentos) > limite:
        _ppt_parrafo(
            slide, f'… y {len(departamentos) - limite} departamento(s) adicional(es) — ver detalle en el informe Word/PDF.',
            top=top + PptxInches(0.4 * n_filas) + PptxInches(0.1), tamano=10, height=PptxInches(0.4),
        )


def _ppt_tabla_formularios_detalle(slide, formularios, top, limite=12):
    """Detalle a nivel de CADA FSC individual — versión slide de
    `_tabla_formularios_detalle_docx`, para el capítulo de subdirección del PPT."""
    filas_datos = formularios[:limite]
    n_filas = len(filas_datos) + 1
    tabla_shape = slide.shapes.add_table(n_filas, 5, PptxInches(PPT_MARGEN), top, PptxInches(PPT_ANCHO_CONTENIDO), PptxInches(0.4 * n_filas))
    tabla = tabla_shape.table
    _ppt_fijar_anchos_columnas(tabla, [1.47, 4.55, 1.47, 2.27, 2.54])
    for i, h in enumerate(['Formulario', 'Unidad Requirente', 'Dentro/Fuera', 'Monto', 'Fecha Derivado']):
        tabla.cell(0, i).text = h
    for r, f in enumerate(filas_datos, start=1):
        tabla.cell(r, 0).text = f['id_formulario'] or f"{f['folio']}/{f['anho']}"
        tabla.cell(r, 1).text = _truncar(f['unidad_requirente'], 45)
        tabla.cell(r, 2).text = 'Dentro' if f['dentro_fuera_pac'] == 'DENTRO' else 'Fuera'
        tabla.cell(r, 3).text = _money(f['monto_estimado'])
        tabla.cell(r, 4).text = f['fecha_derivado'] or '—'
    _ppt_estilizar_tabla(tabla, tamano_fuente=11)
    if len(formularios) > limite:
        _ppt_parrafo(
            slide, f'… y {len(formularios) - limite} formulario(s) adicional(es) — ver detalle en el informe Word/PDF.',
            top=top + PptxInches(0.4 * n_filas) + PptxInches(0.1), tamano=10, height=PptxInches(0.4),
        )


def _ppt_tabla_fichas_ejecucion(slide, fichas, top, limite=10):
    """Detalle de proyectos del Plan de Compras con CON QUÉ formulario(s) se
    ejecutó cada uno — versión slide de `_tabla_fichas_ejecucion_detalle_docx`."""
    filas_datos = fichas[:limite]
    n_filas = len(filas_datos) + 1
    tabla_shape = slide.shapes.add_table(n_filas, 4, PptxInches(PPT_MARGEN), top, PptxInches(PPT_ANCHO_CONTENIDO), PptxInches(0.4 * n_filas))
    tabla = tabla_shape.table
    _ppt_fijar_anchos_columnas(tabla, [3.08, 2.54, 1.47, 5.21])
    for i, h in enumerate(['Proyecto', 'Departamento', 'Estado', 'Ejecutado con (Formulario)']):
        tabla.cell(0, i).text = h
    for r, f in enumerate(filas_datos, start=1):
        ejecutores = f.get('formularios_ejecutores') or []
        texto_ejecutores = ', '.join(fx['id_formulario'] or f"{fx['folio']}/{fx['anho']}" for fx in ejecutores) if ejecutores else '—'
        tabla.cell(r, 0).text = _truncar(f['nombre_proyecto'] or f['id_proyecto'], 38)
        tabla.cell(r, 1).text = _truncar(f['depto_nombre'] or f['depto_texto'], 30)
        tabla.cell(r, 2).text = f['estado_ejecucion']
        tabla.cell(r, 3).text = _truncar(texto_ejecutores, 45)
    _ppt_estilizar_tabla(tabla, tamano_fuente=11)
    if len(fichas) > limite:
        _ppt_parrafo(
            slide, f'… y {len(fichas) - limite} proyecto(s) adicional(es) — ver detalle en el informe Word/PDF.',
            top=top + PptxInches(0.4 * n_filas) + PptxInches(0.1), tamano=10, height=PptxInches(0.4),
        )


def _ppt_tabla_fichas_pac(slide, fichas, top, limite=10):
    """Fichas del Plan de Compras (próximo mes/atrasadas) — versión slide de
    `_tabla_fichas_pac_docx`, usada por Alertas y Seguimiento agrupado por subdirección."""
    filas_datos = fichas[:limite]
    n_filas = len(filas_datos) + 1
    tabla_shape = slide.shapes.add_table(n_filas, 4, PptxInches(PPT_MARGEN), top, PptxInches(PPT_ANCHO_CONTENIDO), PptxInches(0.4 * n_filas))
    tabla = tabla_shape.table
    _ppt_fijar_anchos_columnas(tabla, [4.55, 3.21, 2.27, 2.27])
    for i, h in enumerate(['Proyecto', 'Departamento', 'Fecha Compra', 'Estado']):
        tabla.cell(0, i).text = h
    for r, f in enumerate(filas_datos, start=1):
        tabla.cell(r, 0).text = _truncar(f.get('nombre_proyecto') or f['id_proyecto'], 42)
        tabla.cell(r, 1).text = _truncar(f['depto_nombre'] or f['depto_texto'], 32)
        tabla.cell(r, 2).text = f['fecha_mas_proxima'] or '—'
        tabla.cell(r, 3).text = f['estado_ejecucion']
    _ppt_estilizar_tabla(tabla, tamano_fuente=11)
    if len(fichas) > limite:
        _ppt_parrafo(
            slide, f'… y {len(fichas) - limite} ficha(s) adicional(es) — ver detalle en el informe Word/PDF.',
            top=top + PptxInches(0.4 * n_filas) + PptxInches(0.1), tamano=10, height=PptxInches(0.4),
        )


def _ppt_tabla_proyectos_sin_iniciar(slide, proyectos, top, limite=10):
    filas_datos = proyectos[:limite]
    n_filas = len(filas_datos) + 1
    tabla_shape = slide.shapes.add_table(n_filas, 3, PptxInches(PPT_MARGEN), top, PptxInches(PPT_ANCHO_CONTENIDO), PptxInches(0.4 * n_filas))
    tabla = tabla_shape.table
    _ppt_fijar_anchos_columnas(tabla, [6.69, 2.94, 2.67])
    for i, h in enumerate(['Proyecto', 'Fecha Planificada', 'Estado']):
        tabla.cell(0, i).text = h
    for r, p in enumerate(filas_datos, start=1):
        tabla.cell(r, 0).text = _truncar(p.get('nombre_proyecto') or p['id_proyecto'], 55)
        tabla.cell(r, 1).text = p['fecha_inicio_compra']
        tabla.cell(r, 2).text = 'Pendiente' if p['estado'] == 'PENDIENTE' else 'Atrasado'
    _ppt_estilizar_tabla(tabla, tamano_fuente=11)
    if len(proyectos) > limite:
        _ppt_parrafo(
            slide, f'… y {len(proyectos) - limite} proyecto(s) adicional(es) — ver detalle en el informe Word/PDF.',
            top=top + PptxInches(0.4 * n_filas) + PptxInches(0.1), tamano=10, height=PptxInches(0.4),
        )


def _ppt_capitulo_subdireccion(prs, nombre_display, fsc, planer, responsables, anho, formularios_detalle=None, fichas_detalle=None):
    """Mini-capítulo de la subdirección `nombre_display` en el PPT — decisión del
    usuario 2026-07-21 (2ª ronda): expandir el PPT de 'ejecutivo resumido' a una
    presentación completa organizable por subdirección, para exponer a cada una
    individualmente. Espejo del capítulo de Word/PDF pero en formato slide: KPIs +
    donut, tabla+gráfico Dentro/Fuera por departamento, tabla+gráfico de Ejecución
    del Plan de Compras por departamento, y (2026-07-23) detalle a nivel de cada FSC
    individual y de cada proyecto/ficha con sus formularios ejecutores. `fsc`/`planer`
    pueden ser `None` si esa subdirección no tiene datos en el dominio correspondiente
    (períodos cortos)."""
    formularios_detalle = formularios_detalle or []
    fichas_detalle = fichas_detalle or []
    # --- Portada de capítulo: franja + tarjetas KPI + donuts -----------------
    slide = _ppt_slide_en_blanco(prs)
    franja = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0.08), PptxInches(PPT_ANCHO), PptxInches(1.1))
    franja.fill.solid()
    franja.fill.fore_color.rgb = COLOR_TITULO_PPT
    franja.line.fill.background()
    franja.shadow.inherit = False
    tf = franja.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = nombre_display
    run.font.size = PptxPt(24)
    run.font.bold = True
    run.font.color.rgb = PptxRGBColor(0xff, 0xff, 0xff)

    n_tiles = 4
    gap = 0.15
    ancho_tile = (PPT_ANCHO_CONTENIDO - (n_tiles - 1) * gap) / n_tiles
    tiles = []
    if fsc:
        tiles.append((_n(fsc['total']), 'Formularios Totales', COLOR_TITULO_PPT))
        tiles.append((f"{fsc['pct_dentro']}%", 'Dentro del PAC', COLOR_PPT_VERDE if fsc['pct_dentro'] >= 50 else COLOR_PPT_AMBAR))
    else:
        tiles.append(('—', 'Formularios Totales', COLOR_TITULO_PPT))
        tiles.append(('—', 'Dentro del PAC', COLOR_TITULO_PPT))
    if planer:
        tiles.append((_n(planer['total']), f'Fichas Plan {anho}', COLOR_PPT_AZUL_PROFUNDO))
        tiles.append((f"{planer['pct_ejecutado']}%", '% Ejecutado', COLOR_PPT_VERDE if planer['pct_ejecutado'] >= 70 else COLOR_PPT_AMBAR))
    else:
        tiles.append(('—', f'Plan de Compras {anho}', COLOR_PPT_AZUL_PROFUNDO))
        tiles.append(('—', '% Ejecutado', COLOR_PPT_AZUL_PROFUNDO))
    for i, (valor, etiqueta, color) in enumerate(tiles):
        left = PptxInches(PPT_MARGEN + i * (ancho_tile + gap))
        _ppt_tarjeta_kpi(slide, left, PptxInches(1.35), PptxInches(ancho_tile), PptxInches(0.85), valor, etiqueta, color_bg=color)

    top_donas = PptxInches(2.4)
    if responsables:
        _ppt_parrafo(
            slide, 'Responsables: ' + ', '.join(responsables[:5]) + ('…' if len(responsables) > 5 else ''),
            top=PptxInches(2.35), tamano=10.5, height=PptxInches(0.3),
        )
        top_donas = PptxInches(2.8)
    if fsc:
        _ppt_imagen(slide, grafico_donut_dentro_fuera(fsc['pct_dentro'], titulo='Dentro / Fuera PAC'),
                    PptxInches(3.17), top_donas, PptxInches(3.2))
    if planer:
        _ppt_imagen(
            slide,
            grafico_donut_ejecucion_planer(planer['ejecutados'], planer['pendientes'], planer['atrasados']),
            PptxInches(6.97), top_donas, PptxInches(3.2),
        )

    # --- Dentro/Fuera por departamento: tabla + gráfico ----------------------
    if fsc and fsc['departamentos']:
        slide = _ppt_slide_en_blanco(prs)
        _ppt_titulo(slide, f'{nombre_display} — Formularios por Departamento')
        _ppt_tabla_depto_dentro_fuera(slide, fsc['departamentos'], top=PptxInches(1.1))

        img_ranking = grafico_barras_ranking(
            fsc['departamentos'], campo_nombre='nombre', campo_valor='pct_dentro',
            titulo='Ranking de departamentos — % Dentro PAC', color=COLOR_DENTRO,
        )
        if img_ranking:
            slide = _ppt_slide_en_blanco(prs)
            _ppt_titulo(slide, f'{nombre_display} — Ranking % Dentro PAC')
            _ppt_imagen(slide, img_ranking, PptxInches(1.87), PptxInches(1.1), PptxInches(9.6))

    # --- Detalle de formularios individuales (2026-07-23) --------------------
    if formularios_detalle:
        slide = _ppt_slide_en_blanco(prs)
        _ppt_titulo(slide, f'{nombre_display} — Detalle de Formularios')
        _ppt_parrafo(
            slide,
            f'Detalle de los {_n(len(formularios_detalle))} formularios individuales de {nombre_display} en el período.',
            top=PptxInches(1.0), tamano=11, height=PptxInches(0.4),
        )
        _ppt_tabla_formularios_detalle(slide, formularios_detalle, top=PptxInches(1.45))

    # --- Ejecución PAC por departamento: tabla + gráfico ---------------------
    if planer and planer['departamentos']:
        slide = _ppt_slide_en_blanco(prs)
        _ppt_titulo(slide, f'{nombre_display} — Ejecución del Plan de Compras por Departamento')
        _ppt_tabla_depto_ejecucion(slide, planer['departamentos'], top=PptxInches(1.1))

        img_ejec_depto = grafico_barras_ejecucion_por_depto(
            planer['departamentos'], titulo=f'{nombre_display} — Ejecutado / Pendiente / Atrasado por depto.',
        )
        if img_ejec_depto:
            slide = _ppt_slide_en_blanco(prs)
            _ppt_titulo(slide, f'{nombre_display} — Ejecución por Departamento')
            _ppt_imagen(slide, img_ejec_depto, PptxInches(1.87), PptxInches(1.1), PptxInches(9.6))

    # --- Detalle de proyectos y con qué formulario se ejecutaron (2026-07-23) ----
    if fichas_detalle:
        slide = _ppt_slide_en_blanco(prs)
        _ppt_titulo(slide, f'{nombre_display} — Proyectos del Plan de Compras')
        _ppt_parrafo(
            slide,
            f'Detalle de los {_n(len(fichas_detalle))} proyectos del Plan de Compras {anho} de {nombre_display}, '
            'indicando con qué formulario(s) se ejecutó cada uno.',
            top=PptxInches(1.0), tamano=11, height=PptxInches(0.5),
        )
        _ppt_tabla_fichas_ejecucion(slide, fichas_detalle, top=PptxInches(1.55))


def generar_presentacion_ppt(periodo):
    """Genera la presentación PPT institucional del período ('YYYY-MM' o 'YYYY-QN').
    Retorna BytesIO. Estructura: sección general (KPIs institucionales, comparativas,
    evolución mensual, resumen por subdirección) seguida de UN capítulo completo por
    cada subdirección institucional (KPIs propios, gráficos y tablas por departamento
    — ver `_ppt_capitulo_subdireccion`), y cierra con rankings/alertas/conclusiones.
    Funciona igual para período mensual ('YYYY-MM') o trimestral ('YYYY-QN') — ambos
    pasan por `_datos_informe_completo`, que ya resuelve el rango de fechas."""
    d = _datos_informe_completo(periodo)
    label, anho, anho_ant, hoy = d['label'], d['anho'], d['anho_anterior'], d['hoy']

    prs = Presentation()
    prs.slide_width = PptxInches(PPT_ANCHO)
    prs.slide_height = PptxInches(PPT_ALTO)

    # --- Portada -------------------------------------------------------------
    _ppt_portada(
        prs,
        badge_texto='INFORME INSTITUCIONAL',
        titulo_lineas=['Informe de Seguimiento y Ejecución', 'del Plan Anual de Compras'],
        subtitulo_inst=NOMBRE_INSTITUCION.upper(),
        meta_texto=f'Período: {label} · Año PAC {anho} · Generado el {hoy.strftime("%d-%m-%Y")}',
        logo_bytes=_logo_bytes(), imagen_bytes=_edificio_bytes(),
    )

    # --- Agenda ----------------------------------------------------------------
    slide = _ppt_slide_en_blanco(prs)
    _ppt_titulo(slide, 'Agenda')
    subdirecciones_institucionales = [s for s in d['subdirecciones'] if s['institucional']]
    agenda = [
        'Resumen Ejecutivo Institucional',
        'Comparativa Histórica y Evolución Mensual',
        *([f'Avance Trimestral — {label}'] if d['avance_trimestral'] else []),
        'Ejecución del Plan Anual de Compras',
        'Resumen General por Subdirección',
        *[f'Detalle — {_nombre_subdireccion_display(s["nombre"])}' for s in subdirecciones_institucionales],
        'Rankings Institucionales',
        'Alertas y Seguimiento',
        'Conclusiones y Recomendaciones',
    ]
    n_agenda = len(agenda)
    gap_agenda = 0.08
    alto_disponible_agenda = 6.0
    alto_tile_agenda = max(0.4, min(0.58, (alto_disponible_agenda - gap_agenda * (n_agenda - 1)) / n_agenda))
    top_agenda = PptxInches(1.0)
    for i, item in enumerate(agenda):
        _ppt_tarjeta_numerada(slide, top_agenda, f'{i + 1:02d}', item, height=alto_tile_agenda)
        top_agenda += PptxInches(alto_tile_agenda + gap_agenda)

    # --- Separador — Resumen Ejecutivo ---------------------------------------
    _ppt_divisor_seccion(prs, '01', ['RESUMEN EJECUTIVO', 'INSTITUCIONAL'], f'Período: {label} · Año PAC {anho}')

    # --- Resumen Ejecutivo -------------------------------------------------
    slide = _ppt_slide_en_blanco(prs)
    _ppt_titulo(slide, 'Resumen Ejecutivo Institucional', subtitulo=f'Período: {label}')
    kpis = d['dentro_fuera']['kpis']
    var_ant = d['comparativa_periodos']['periodo_anterior']['variacion_pp']
    top_fila = PptxInches(1.0)
    _ppt_fila_indicador(
        slide, top_fila, 'Formularios Dentro del PAC', f"{kpis['pct_dentro']}%",
        COLOR_PPT_VERDE if kpis['pct_dentro'] >= 50 else COLOR_PPT_AMBAR,
        comparacion=f"Total: {_n(kpis['total'])}  ·  Dentro: {_n(kpis['dentro'])}  ·  Fuera: {_n(kpis['fuera'])}",
        explicacion=f"Variación vs período anterior: {var_ant:+.1f} p.p.",
        flecha='↑' if var_ant >= 0 else '↓',
    )
    top_fila += PptxInches(0.76)
    pct_en_fecha = d['temporal']['kpis']['pct_en_fecha']
    _ppt_fila_indicador(
        slide, top_fila, 'Cumplimiento Temporal', f"{pct_en_fecha}%",
        COLOR_PPT_VERDE if pct_en_fecha >= 70 else COLOR_PPT_AMBAR,
        explicacion='% de fichas del Plan de Compras ejecutadas dentro del plazo planificado.', flecha='→',
    )
    top_fila += PptxInches(0.76)
    _ppt_fila_indicador(
        slide, top_fila, f'Ejecución Plan de Compras {anho}', f"{d['pct_ejecutado_fichas']}%",
        COLOR_PPT_VERDE if d['pct_ejecutado_fichas'] >= 70 else COLOR_PPT_AMBAR,
        comparacion=f"{_n(d['total_fichas'])} fichas totales",
        explicacion=f"{_n(d['ejecutados_fichas'])} ejecutadas con formulario u OC enlazada.", flecha='→',
    )
    _ppt_imagen(slide, grafico_donut_dentro_fuera(kpis['pct_dentro']), PptxInches(3.97), PptxInches(3.5), PptxInches(2.4))
    _ppt_imagen(slide, grafico_donut_cumplimiento_temporal(d['temporal']['kpis']), PptxInches(6.97), PptxInches(3.5), PptxInches(2.4))

    # --- Comparativa histórica + Evolución mensual --------------------------
    img_comp = grafico_barras_comparativa_anual(d['dentro_fuera']['comparativa_anual'])
    if img_comp:
        slide = _ppt_slide_en_blanco(prs)
        _ppt_titulo(slide, 'Comparativa Histórica — Dentro vs Fuera PAC')
        _ppt_imagen(slide, img_comp, PptxInches(2.17), PptxInches(1.2), PptxInches(9))

    img_evol = grafico_evolucion_mensual_dentro_fuera(d['temporalidad_mensual']['meses'], anho, anho_ant)
    if img_evol:
        slide = _ppt_slide_en_blanco(prs)
        _ppt_titulo(slide, f'Evolución Mensual {anho} vs {anho_ant}')
        _ppt_imagen(slide, img_evol, PptxInches(1.4), PptxInches(1.2), PptxInches(10.5))

    # --- Avance Trimestral (solo si el período es un trimestre) --------------
    if d['avance_trimestral']:
        av = d['avance_trimestral']
        slide = _ppt_slide_en_blanco(prs)
        _ppt_titulo(slide, f'Avance Trimestral — {label}')
        texto_var_dentro = (
            f'{av["variacion_pct_dentro"]:+.1f} p.p.' if av['variacion_pct_dentro'] is not None else 'sin datos suficientes'
        )
        texto_var_ejec = (
            f'{av["variacion_pct_ejecutado"]:+.1f} p.p.' if av['variacion_pct_ejecutado'] is not None else 'sin datos suficientes'
        )
        _ppt_parrafo(
            slide,
            f'Variación % Dentro PAC dentro del trimestre: {texto_var_dentro}   ·   '
            f'Variación % Ejecutado del Plan de Compras: {texto_var_ejec}',
            top=PptxInches(1.05), tamano=13, height=PptxInches(0.6),
        )
        img_av_fsc = grafico_evolucion_mensual_dentro_fuera(av['meses_fsc'], anho, anho_ant, titulo=f'% Dentro PAC — {label}')
        _ppt_imagen(slide, img_av_fsc, PptxInches(0.5), PptxInches(1.8), PptxInches(6.0))
        img_av_planer = grafico_barras_ejecucion_mensual_planer(av['meses_planer'], titulo=f'Ejecución — {label}')
        _ppt_imagen(slide, img_av_planer, PptxInches(6.8), PptxInches(1.8), PptxInches(6.0))

    # --- Ejecución del Plan de Compras ---------------------------------------
    slide = _ppt_slide_en_blanco(prs)
    _ppt_titulo(slide, 'Ejecución del Plan Anual de Compras')
    _ppt_parrafo(
        slide,
        f'Año PAC {anho}: {_n(d["total_fichas"])} fichas — {d["pct_ejecutado_fichas"]}% ejecutadas '
        f'({_n(d["ejecutados_fichas"])} con formulario u OC enlazada). '
        f'Monto total planificado: {_money(d["monto_total_fichas"])}.',
        top=PptxInches(1.05), tamano=14, height=PptxInches(1.0),
    )
    img_ejec = grafico_barras_ejecucion_mensual_planer(d['temporal_mensual_planer']['meses'])
    if img_ejec:
        _ppt_imagen(slide, img_ejec, PptxInches(1.5), PptxInches(2.1), PptxInches(10.3))

    # --- Resumen general por subdirección (1 tabla combinada, panorama antes
    # del detalle) ------------------------------------------------------------
    slide = _ppt_slide_en_blanco(prs)
    _ppt_titulo(slide, 'Resumen General por Subdirección')
    _ppt_parrafo(
        slide,
        'Formularios Dentro/Fuera del PAC y ejecución del Plan de Compras, por subdirección institucional — '
        'el detalle completo de cada una (KPIs, gráficos y departamentos) se desarrolla a continuación.',
        top=PptxInches(1.0), tamano=12, height=PptxInches(0.7),
    )
    _ppt_tabla_resumen_combinado(slide, d['subdirecciones'], top=PptxInches(1.75))

    # --- Separador — Detalle por Subdirección + capítulos --------------------
    _ppt_divisor_seccion(prs, '02', ['DETALLE POR', 'SUBDIRECCIÓN'], 'KPIs, gráficos y departamentos de cada subdirección institucional')
    for s in subdirecciones_institucionales:
        _ppt_capitulo_subdireccion(
            prs, _nombre_subdireccion_display(s['nombre']), s['fsc'], s['planer'], s['responsables'], anho,
            formularios_detalle=s['formularios_detalle'], fichas_detalle=s['fichas_detalle'],
        )

    # --- Separador — Rankings -------------------------------------------------
    _ppt_divisor_seccion(prs, '03', ['RANKINGS', 'INSTITUCIONALES'], 'Mejores departamentos y mayor oportunidad de mejora')

    # --- Rankings ----------------------------------------------------------
    slide = _ppt_slide_en_blanco(prs)
    _ppt_titulo(slide, 'Ranking Institucional — Mejores Departamentos')
    if d['rankings_depto']['mejores']:
        _ppt_tabla_ranking(slide, d['rankings_depto']['mejores'], top=PptxInches(1.2), limite=5)
    else:
        _ppt_parrafo(slide, 'Sin datos suficientes para generar el ranking en este período.', top=PptxInches(1.2))

    slide = _ppt_slide_en_blanco(prs)
    _ppt_titulo(slide, 'Ranking Institucional — Mayor Oportunidad de Mejora')
    if d['rankings_depto']['peores']:
        _ppt_tabla_ranking(slide, d['rankings_depto']['peores'], top=PptxInches(1.2), limite=5)
    else:
        _ppt_parrafo(slide, 'Sin datos suficientes para generar el ranking en este período.', top=PptxInches(1.2))

    # --- Ranking por Formulario (ausente hasta la revisión de código 2026-07-27 —
    # Word/PDF ya incluían esta dimensión, el PPT solo mostraba departamentos) -----
    slide = _ppt_slide_en_blanco(prs)
    _ppt_titulo(slide, 'Ranking de Formularios — Mejor Desempeño')
    if d['rankings_formulario']['mejores']:
        _ppt_tabla_ranking_formulario(slide, d['rankings_formulario']['mejores'], top=PptxInches(1.2), limite=5)
    else:
        _ppt_parrafo(slide, 'Sin datos suficientes para generar el ranking en este período.', top=PptxInches(1.2))

    slide = _ppt_slide_en_blanco(prs)
    _ppt_titulo(slide, 'Ranking de Formularios — Mayor Oportunidad de Mejora')
    if d['rankings_formulario']['peores']:
        _ppt_tabla_ranking_formulario(slide, d['rankings_formulario']['peores'], top=PptxInches(1.2), limite=5)
    else:
        _ppt_parrafo(slide, 'Sin datos suficientes para generar el ranking en este período.', top=PptxInches(1.2))

    # --- Separador — Alertas y Seguimiento ------------------------------------
    _ppt_divisor_seccion(prs, '04', ['ALERTAS Y', 'SEGUIMIENTO'], 'Proyectos sin iniciar, fichas próximas y atrasadas')

    # --- Alertas y Seguimiento: resumen institucional + detalle por subdirección
    # (2026-07-23, para mayor claridad) --------------------------------------
    n_sin_iniciar = _total_items_grupos(d['proyectos_sin_iniciar'])
    n_proximo_mes = _total_items_grupos(d['proximo_mes'])
    n_atrasadas = _total_items_grupos(d['atrasadas'])
    slide = _ppt_slide_en_blanco(prs)
    _ppt_titulo(slide, 'Alertas y Seguimiento')
    _ppt_tarjeta_numerada(
        slide, PptxInches(1.2), '⚠', f'Proyectos planificados sin formulario Dentro PAC derivado: {_n(n_sin_iniciar)}',
        color=COLOR_PPT_ROJO, height=0.75,
    )
    _ppt_tarjeta_numerada(
        slide, PptxInches(2.15), '📅', f'Fichas a ejecutar el próximo mes (Año PAC {anho}): {_n(n_proximo_mes)}',
        color=COLOR_PPT_AZUL_PROFUNDO, height=0.75,
    )
    _ppt_tarjeta_numerada(
        slide, PptxInches(3.1), '⏰', f'Fichas atrasadas (Año PAC {anho}): {_n(n_atrasadas)}',
        color=COLOR_PPT_AMBAR, height=0.75,
    )
    _ppt_parrafo(slide, 'Detalle por subdirección en las diapositivas siguientes.', top=PptxInches(4.1), tamano=12, height=PptxInches(0.4))

    for grupo in d['proyectos_sin_iniciar']:
        slide = _ppt_slide_en_blanco(prs)
        _ppt_titulo(slide, f'Proyectos sin iniciar — {grupo["nombre_display"]}')
        _ppt_tabla_proyectos_sin_iniciar(slide, grupo['items'], top=PptxInches(1.1))

    for grupo in d['proximo_mes']:
        slide = _ppt_slide_en_blanco(prs)
        _ppt_titulo(slide, f'Fichas a ejecutar próximo mes — {grupo["nombre_display"]}')
        _ppt_tabla_fichas_pac(slide, grupo['items'], top=PptxInches(1.1))

    for grupo in d['atrasadas']:
        slide = _ppt_slide_en_blanco(prs)
        _ppt_titulo(slide, f'Fichas atrasadas — {grupo["nombre_display"]}')
        _ppt_tabla_fichas_pac(slide, grupo['items'], top=PptxInches(1.1))

    # --- Separador — Conclusiones ---------------------------------------------
    _ppt_divisor_seccion(prs, '05', ['CONCLUSIONES Y', 'RECOMENDACIONES'], 'Cierre y próximos pasos del Plan Anual de Compras')

    # --- Conclusiones ----------------------------------------------------------
    slide = _ppt_slide_en_blanco(prs)
    _ppt_titulo(slide, 'Conclusiones y Recomendaciones')
    _ppt_parrafo(
        slide, parrafo_conclusiones(label, d['dentro_fuera']['kpis'], d['temporal']['kpis']),
        top=PptxInches(1.2), tamano=15, height=PptxInches(2.6),
    )
    _ppt_parrafo(
        slide,
        f'Ejecución del Plan de Compras {anho}: {d["pct_ejecutado_fichas"]}% ejecutado. '
        + ('Se recomienda priorizar la regularización de fichas atrasadas.' if d['pct_ejecutado_fichas'] < 70
           else 'Nivel de ejecución satisfactorio — mantener el seguimiento mensual.'),
        top=PptxInches(4.2), tamano=15, height=PptxInches(1.4),
    )

    for slide in prs.slides:
        _ppt_pie_pagina(slide)
    _ppt_numerar_diapositivas(prs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
