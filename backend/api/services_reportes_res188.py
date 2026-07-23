"""
Reportería institucional (Word, PDF, PPT) para los Indicadores Res.188/2026 del
módulo PAC (`/pac`). Empieza cubriendo solo el Indicador 1 — % Compras dentro
del PAC — con el detalle de enlace PAC (evolución mensual, trimestral,
comparativo anual, corregidas, etc.) ya disponible en `calcular_oc_stats`.

Arquitectura pensada para escalar: a diferencia de `services_reportes.py`
(3 funciones monolíticas que repiten la secuencia de secciones a mano por
formato), este módulo arma primero una lista de `Seccion` (contenido puro,
sin conocimiento de docx/pptx/reportlab) vía `_construir_secciones()`, y cada
generador (`generar_informe_word_ind1`, `generar_reporte_pdf_ind1`,
`generar_presentacion_ppt_ind1`) es un solo loop que traduce esa MISMA lista
a su formato. Agregar el Indicador 2 (o cualquier sección nueva) más adelante
significa agregar `Seccion`es en `_construir_secciones()` — no tocar los 3
renderers.

Reutiliza deliberadamente los helpers de bajo nivel ya probados en
`services_reportes.py` (colores institucionales, `_fig_a_bytes`, márgenes,
títulos de capítulo/sección, estilos PDF, helpers de slide PPT) para que el
resultado visual sea indistinguible del informe de Cumplimiento Interno PAC.
Lo único propio de este módulo es la portada/pie de página (llevan un título
distinto) y, por supuesto, el contenido.
"""
import io
from dataclasses import dataclass, field
from datetime import date

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.platypus import BaseDocTemplate, PageTemplate, Frame, Paragraph, Spacer, Table, PageBreak
from reportlab.platypus.tableofcontents import TableOfContents
from reportlab.lib.styles import ParagraphStyle

from .services import calcular_indicadores_res188, calcular_oc_stats
from .plantillas_narrativas import _n, _money
from .services_reportes import (
    NOMBRE_INSTITUCION,
    COLOR_INSTITUCIONAL, COLOR_INSTITUCIONAL_CLARO, COLOR_DENTRO, COLOR_FUERA, COLOR_PENDIENTE,
    _fig_a_bytes, _logo_bytes, _edificio_bytes,
    _configurar_margenes_docx, _agregar_campo_docx, _forzar_actualizacion_campos_docx,
    _titulo_capitulo_docx, _titulo_seccion_docx,
    _PDF_ESTILOS, _PDF_TABLA_ESTILO, _celda, _celda_encabezado, _fila_encabezado, _pdf_imagen,
    _ppt_slide_en_blanco, _ppt_titulo, _ppt_parrafo, _ppt_imagen, _ppt_pie_pagina, _ppt_numerar_diapositivas,
    COLOR_TITULO_PPT, COLOR_TEXTO_PPT,
)

import matplotlib.pyplot as plt  # backend Agg ya fijado por el import de services_reportes arriba

TITULO_INFORME_IND1 = 'Informe de Gestión PAC — Indicador 1: % Compras dentro del PAC'

MESES_CORTOS = ['Ene', 'Feb', 'Mar', 'Abr', 'May', 'Jun', 'Jul', 'Ago', 'Sep', 'Oct', 'Nov', 'Dic']
TRIMESTRES = [('T1 (Ene-Mar)', (1, 2, 3)), ('T2 (Abr-Jun)', (4, 5, 6)), ('T3 (Jul-Sep)', (7, 8, 9)), ('T4 (Oct-Dic)', (10, 11, 12))]
PALETA_ANIOS = ['#1e3a5f', '#38b2bd', '#16a34a', '#d97706', '#dc2626', '#7c3aed']


def _fmt_monto_corto(x):
    """$1.234.567.890 -> '$1.2MM' ; $45.000.000 -> '$45M' — para ejes de gráfico."""
    if abs(x) >= 1e9:
        return f'${x / 1e9:.1f}MM'
    if abs(x) >= 1e6:
        return f'${x / 1e6:.0f}M'
    return _money(x)


# =============================================================================
# Capa de datos — une calcular_indicadores_res188 + calcular_oc_stats
# =============================================================================

def _construir_datos(anio):
    return {
        'anio': anio,
        'hoy': date.today(),
        'indicadores': calcular_indicadores_res188(anio),
        'oc': calcular_oc_stats(anio),
    }


# =============================================================================
# Modelo de contenido — independiente de docx/pptx/reportlab
# =============================================================================

@dataclass
class Grafico:
    imagen: object                 # BytesIO (PNG) o None
    caption: str
    ancho_word: float = 5.6        # inches
    ancho_pdf: tuple = (5.6, 3.15)  # (w, h) inches
    ancho_ppt: float = 8.6         # inches


@dataclass
class Tabla:
    encabezados: list
    filas: list                    # list[list[str]]
    anchos_word: list = None       # inches por columna (None = reparto automático)
    anchos_pdf: list = None


@dataclass
class Seccion:
    titulo: str
    nivel: int = 1                  # 1 = capítulo, 2 = subsección
    parrafo: str = None
    graficos: list = field(default_factory=list)   # list[Grafico]
    tabla: Tabla = None
    nota: str = None


# =============================================================================
# Gráficos (matplotlib, paleta institucional, devuelven BytesIO vía _fig_a_bytes)
# =============================================================================

def _grafico_donut_enlace(pct_enlace):
    if pct_enlace is None:
        return None
    pct_fuera = round(100 - pct_enlace, 1)
    fig, ax = plt.subplots(figsize=(3.6, 3.6))
    ax.pie(
        [pct_enlace, pct_fuera], colors=[COLOR_DENTRO, COLOR_FUERA], startangle=90,
        wedgeprops={'width': 0.35, 'edgecolor': 'white', 'linewidth': 2},
        autopct=lambda p: f'{p:.0f}%' if p > 5 else '', pctdistance=0.82,
        textprops={'color': 'white', 'fontweight': 'bold', 'fontsize': 11},
    )
    ax.set_title('% Enlace al PAC', fontsize=11, pad=10)
    fig.legend(['Enlazada', 'No Enlazada'], loc='lower center', ncol=2, frameon=False,
               bbox_to_anchor=(0.5, -0.02), fontsize=9)
    return _fig_a_bytes(fig)


def _grafico_barras_enlace(labels, enlazada, no_enlazada, titulo):
    fig, ax = plt.subplots(figsize=(7.0, 3.4))
    x = range(len(labels))
    ancho = 0.38
    ax.bar([i - ancho / 2 for i in x], enlazada, width=ancho, label='Enlazada', color=COLOR_DENTRO, edgecolor='none')
    ax.bar([i + ancho / 2 for i in x], no_enlazada, width=ancho, label='No Enlazada', color=COLOR_FUERA, edgecolor='none')
    ax.set_xticks(list(x))
    ax.set_xticklabels(labels, fontsize=9)
    ax.yaxis.set_major_formatter(lambda v, _: _fmt_monto_corto(v))
    ax.set_title(titulo, fontsize=11, pad=10)
    ax.legend(frameon=False, fontsize=9, loc='upper left')
    ax.spines[['top', 'right']].set_visible(False)
    fig.tight_layout()
    return _fig_a_bytes(fig)


def _grafico_evolucion_mensual(evolucion_enlace, anio):
    mapa = {r['mes']: r for r in (evolucion_enlace or [])}
    enl = [mapa.get(m, {}).get('enlazada', 0) for m in range(1, 13)]
    noenl = [mapa.get(m, {}).get('no_enlazada', 0) for m in range(1, 13)]
    return _grafico_barras_enlace(MESES_CORTOS, enl, noenl, f'Evolución Mensual: Enlazada vs No Enlazada — Monto ({anio})')


def _grafico_barras_enlace_cantidad(labels, enlazada, no_enlazada, titulo):
    """Igual que `_grafico_barras_enlace` pero con eje Y entero (N° OC) en vez de monto."""
    fig, ax = plt.subplots(figsize=(7.0, 3.4))
    x = range(len(labels))
    ancho = 0.38
    ax.bar([i - ancho / 2 for i in x], enlazada, width=ancho, label='Enlazada', color=COLOR_DENTRO, edgecolor='none')
    ax.bar([i + ancho / 2 for i in x], no_enlazada, width=ancho, label='No Enlazada', color=COLOR_FUERA, edgecolor='none')
    ax.set_xticks(list(x))
    ax.set_xticklabels(labels, fontsize=9)
    ax.set_ylabel('N° de OC', fontsize=9)
    ax.set_title(titulo, fontsize=11, pad=10)
    ax.legend(frameon=False, fontsize=9, loc='upper left')
    ax.spines[['top', 'right']].set_visible(False)
    fig.tight_layout()
    return _fig_a_bytes(fig)


def _grafico_evolucion_mensual_cantidad(evolucion_enlace, anio):
    mapa = {r['mes']: r for r in (evolucion_enlace or [])}
    enl = [mapa.get(m, {}).get('cantidad_enlazada', 0) for m in range(1, 13)]
    noenl = [mapa.get(m, {}).get('cantidad_no_enlazada', 0) for m in range(1, 13)]
    return _grafico_barras_enlace_cantidad(MESES_CORTOS, enl, noenl, f'Evolución Mensual: Enlazada vs No Enlazada — N° de OC ({anio})')


def _grafico_trimestral(evolucion_enlace, anio):
    mapa = {r['mes']: r for r in (evolucion_enlace or [])}
    labels = [t[0] for t in TRIMESTRES]
    enl = [sum(mapa.get(m, {}).get('enlazada', 0) for m in meses) for _, meses in TRIMESTRES]
    noenl = [sum(mapa.get(m, {}).get('no_enlazada', 0) for m in meses) for _, meses in TRIMESTRES]
    fig, ax = plt.subplots(figsize=(4.6, 3.4))
    x = range(len(labels))
    ancho = 0.38
    ax.bar([i - ancho / 2 for i in x], enl, width=ancho, label='Enlazada', color=COLOR_DENTRO)
    ax.bar([i + ancho / 2 for i in x], noenl, width=ancho, label='No Enlazada', color=COLOR_FUERA)
    ax.set_xticks(list(x))
    ax.set_xticklabels(labels, fontsize=9)
    ax.yaxis.set_major_formatter(lambda v, _: _fmt_monto_corto(v))
    ax.set_title(f'Comparativo Trimestral ({anio})', fontsize=11, pad=10)
    ax.legend(frameon=False, fontsize=9)
    ax.spines[['top', 'right']].set_visible(False)
    fig.tight_layout()
    return _fig_a_bytes(fig)


def _grafico_comparativo_anual(historico_enlace_anual):
    filas = sorted(historico_enlace_anual or [], key=lambda r: r['anio'])
    if not filas:
        return None
    fig, ax = plt.subplots(figsize=(7.0, 3.2))
    anios = [r['anio'] for r in filas]
    pct = [r['pct_enlace'] for r in filas]
    ax.plot(anios, pct, color=COLOR_INSTITUCIONAL, marker='o', linewidth=2.2, markersize=5)
    ax.fill_between(anios, pct, color=COLOR_INSTITUCIONAL, alpha=.12)
    ax.set_ylim(0, 100)
    ax.set_xticks(anios)
    ax.set_ylabel('% Enlace PAC', fontsize=9)
    ax.set_title('Comparativo Anual — % Enlace PAC (histórico institucional)', fontsize=11, pad=10)
    ax.spines[['top', 'right']].set_visible(False)
    for a, p in zip(anios, pct):
        ax.annotate(f'{p:.0f}%', (a, p), textcoords='offset points', xytext=(0, 7), ha='center', fontsize=8, color=COLOR_INSTITUCIONAL)
    fig.tight_layout()
    return _fig_a_bytes(fig)


def _grafico_comparativo_anual_cantidad(historico_enlace_anual):
    filas = sorted(historico_enlace_anual or [], key=lambda r: r['anio'])
    if not filas:
        return None
    labels = [str(r['anio']) for r in filas]
    enl = [r['enlazadas'] for r in filas]
    noenl = [r['total_oc'] - r['enlazadas'] for r in filas]
    return _grafico_barras_enlace_cantidad(labels, enl, noenl, 'Comparativo Anual — N° de OC Enlazadas vs No Enlazadas')


def _grafico_multi_anio_mensual(historico_enlace_mensual, anio_actual):
    por_anio = {}
    for r in (historico_enlace_mensual or []):
        por_anio.setdefault(r['anio'], {})[r['mes']] = r['pct_enlace']
    anios_recientes = sorted(por_anio.keys())[-5:]
    if not anios_recientes:
        return None
    fig, ax = plt.subplots(figsize=(7.2, 3.6))
    for i, a in enumerate(anios_recientes):
        serie = [por_anio.get(a, {}).get(m) for m in range(1, 13)]
        ax.plot(
            range(1, 13), serie, marker='o', markersize=3,
            color=PALETA_ANIOS[i % len(PALETA_ANIOS)],
            linewidth=3 if a == anio_actual else 1.4,
            label=str(a),
        )
    ax.set_xticks(range(1, 13))
    ax.set_xticklabels(MESES_CORTOS, fontsize=9)
    ax.set_ylim(0, 100)
    ax.set_title('% Enlace Mensual — Comparativo con Años Anteriores', fontsize=11, pad=10)
    ax.legend(frameon=False, fontsize=8, ncol=len(anios_recientes))
    ax.spines[['top', 'right']].set_visible(False)
    fig.tight_layout()
    return _fig_a_bytes(fig)


def _grafico_tipo_oc_fuera_pac(no_enlazadas_tipo_oc, limite=8):
    filas = (no_enlazadas_tipo_oc or [])[:limite]
    if not filas:
        return None
    filas = list(reversed(filas))  # barh dibuja de abajo hacia arriba
    fig, ax = plt.subplots(figsize=(6.8, max(2.2, 0.42 * len(filas) + 1)))
    ax.barh([f['tipo_oc'] for f in filas], [f['cantidad'] for f in filas], color=COLOR_FUERA)
    ax.set_title('OC Fuera del PAC por Tipo (N°)', fontsize=11, pad=10)
    ax.spines[['top', 'right']].set_visible(False)
    fig.tight_layout()
    return _fig_a_bytes(fig)


# =============================================================================
# Narrativa (determinística, sin IA — mismo criterio que plantillas_narrativas.py)
# =============================================================================

def _nivel(pct):
    if pct is None:
        return 'sin información suficiente'
    if pct >= 70:
        return 'un nivel satisfactorio'
    if pct >= 40:
        return 'un nivel intermedio, con espacio de mejora'
    return 'un nivel crítico que requiere atención'


def _parrafo_resumen_ejecutivo(datos):
    anio = datos['anio']
    hist = sorted(datos['oc']['historico_enlace_anual'], key=lambda r: r['anio'])
    fila_actual = next((r for r in hist if r['anio'] == anio), None)
    fila_anterior = next((r for r in hist if r['anio'] == anio - 1), None)
    pct = fila_actual['pct_enlace'] if fila_actual else (datos['indicadores'].get('i1') or 0)

    texto = (
        f'Durante {anio}, el {pct:.1f}% del monto total de Órdenes de Compra del Servicio de Salud Osorno '
        f'quedó enlazado al Plan Anual de Compras (PAC), lo que representa {_nivel(pct)}.'
    )
    if fila_anterior and fila_anterior.get('pct_enlace') is not None:
        var = round(pct - fila_anterior['pct_enlace'], 1)
        if var > 1:
            texto += f" Esto es un alza de {var:.1f} puntos porcentuales respecto a {anio - 1} ({fila_anterior['pct_enlace']:.1f}%)."
        elif var < -1:
            texto += f" Esto representa una baja de {abs(var):.1f} puntos porcentuales respecto a {anio - 1} ({fila_anterior['pct_enlace']:.1f}%)."
        else:
            texto += f' Se mantiene relativamente estable respecto a {anio - 1} ({fila_anterior["pct_enlace"]:.1f}%).'
    return texto


def _parrafo_conclusiones(datos):
    anio = datos['anio']
    pct = datos['indicadores'].get('i1') or 0
    corregidas = datos['oc']['corregidas']
    if pct >= 70:
        texto = (
            'El nivel de enlace al Plan Anual de Compras alcanzado es satisfactorio. Se recomienda mantener '
            'las prácticas actuales de planificación y seguir monitoreando mensualmente para sostener la tendencia.'
        )
    else:
        texto = (
            'Se recomienda reforzar la planificación anticipada en las unidades y tipos de compra con menor '
            'apego al PAC, priorizando la revisión de las Órdenes de Compra "Fuera del PAC" con mayor volumen '
            'identificadas en este informe.'
        )
    if corregidas.get('esperando_sync', 0) > 0:
        texto += (
            f" Adicionalmente, existen {_n(corregidas['esperando_sync'])} OC corregidas manualmente que aún no "
            f'han sido confirmadas por la sincronización automática con Mercado Público — se recomienda '
            f'verificarlas en la próxima actualización.'
        )
    return texto


# =============================================================================
# Registro de secciones — el contenido se define UNA vez, los 3 formatos lo consumen igual
# =============================================================================

def _construir_secciones(datos):
    anio = datos['anio']
    ind = datos['indicadores']
    oc = datos['oc']
    corregidas = oc['corregidas']

    secciones = []

    secciones.append(Seccion(
        titulo='Resumen Ejecutivo — Indicador 1',
        parrafo=_parrafo_resumen_ejecutivo(datos),
        graficos=[Grafico(_grafico_donut_enlace(ind.get('i1')), 'Gráfico 1 — Distribución del monto de OC Enlazada vs No Enlazada al PAC.', ancho_word=2.8, ancho_pdf=(2.8, 2.8), ancho_ppt=3.4)],
        tabla=Tabla(
            encabezados=['Indicador', 'Valor'],
            filas=[
                ['Total OC del año', _n(ind.get('total_oc'))],
                ['Monto enlazado al PAC', _money(ind.get('monto_enlazado_pac'))],
                ['% Compras dentro del PAC (Ind.1)', f"{ind.get('i1', 0):.1f}%" if ind.get('i1') is not None else '—'],
            ],
        ),
    ))

    secciones.append(Seccion(
        titulo=f'Evolución Mensual: Enlazada vs No Enlazada — Monto ({anio})',
        graficos=[Grafico(_grafico_evolucion_mensual(oc['evolucion_enlace'], anio), f'Gráfico 2 — Monto mensual de OC Enlazada vs No Enlazada, {anio}.')],
    ))

    secciones.append(Seccion(
        titulo=f'Evolución Mensual: Enlazada vs No Enlazada — N° de OC ({anio})',
        nivel=2,
        parrafo='Mismo período que el gráfico anterior, ahora por cantidad de Órdenes de Compra en vez de monto.',
        graficos=[Grafico(_grafico_evolucion_mensual_cantidad(oc['evolucion_enlace'], anio), f'Gráfico 3 — N° de OC Enlazada vs No Enlazada por mes, {anio}.')],
    ))

    secciones.append(Seccion(
        titulo=f'Comparativo Trimestral ({anio})',
        nivel=2,
        graficos=[Grafico(_grafico_trimestral(oc['evolucion_enlace'], anio), f'Gráfico 4 — Monto trimestral de OC Enlazada vs No Enlazada, {anio}.', ancho_word=4.4, ancho_pdf=(4.4, 3.2), ancho_ppt=6)],
    ))

    hist = sorted(oc['historico_enlace_anual'], key=lambda r: r['anio'])
    secciones.append(Seccion(
        titulo='Comparativo Anual — % Enlace PAC (histórico institucional)',
        parrafo='La siguiente serie muestra la evolución del enlace al PAC en todos los años con Órdenes de Compra registradas, independiente del año seleccionado para el resto del informe.',
        graficos=[Grafico(_grafico_comparativo_anual(hist), 'Gráfico 5 — % de enlace al PAC por año.')],
        tabla=Tabla(
            encabezados=['Año', 'OC Enlazadas', 'Monto Enlazado', 'Monto No Enlazado', '% Enlace'],
            filas=[
                [str(r['anio']), f"{_n(r['enlazadas'])} / {_n(r['total_oc'])}", _money(r['monto_enlazado']), _money(r['monto_no_enlazado']), f"{r['pct_enlace']:.1f}%"]
                for r in hist
            ],
        ),
    ))

    secciones.append(Seccion(
        titulo='Comparativo Anual — N° de OC Enlazadas vs No Enlazadas',
        nivel=2,
        parrafo='Misma serie histórica anterior, ahora por cantidad de Órdenes de Compra en vez de monto.',
        graficos=[Grafico(_grafico_comparativo_anual_cantidad(hist), 'Gráfico 6 — N° de OC Enlazada vs No Enlazada por año.')],
    ))

    secciones.append(Seccion(
        titulo='% Enlace Mensual — Comparativo con Años Anteriores',
        parrafo=f'Serie mensual de los últimos años con datos; la línea de {anio} se resalta para ubicar el año en análisis dentro de la tendencia histórica.',
        graficos=[Grafico(_grafico_multi_anio_mensual(oc['historico_enlace_mensual'], anio), 'Gráfico 7 — % de enlace al PAC por mes, comparado entre años.')],
    ))

    tipo_oc_filas = oc['no_enlazadas_tipo_oc'][:12]
    secciones.append(Seccion(
        titulo=f'Órdenes de Compra Fuera del PAC — Detalle por Tipo ({anio})',
        parrafo='Detalle de las OC que no quedaron enlazadas al PAC durante el período, agrupadas por tipo de compra — base para priorizar la revisión y corrección manual.',
        graficos=[Grafico(_grafico_tipo_oc_fuera_pac(oc['no_enlazadas_tipo_oc']), 'Gráfico 8 — N° de OC fuera del PAC por tipo.', ancho_word=5.2, ancho_pdf=(5.2, 3.4), ancho_ppt=7.5)],
        tabla=Tabla(
            encabezados=['Tipo OC', 'Cantidad', 'Monto Neto'],
            filas=[[f['tipo_oc'], _n(f['cantidad']), _money(f['monto'])] for f in tipo_oc_filas],
        ),
    ))

    matriz = oc['matriz_tipo_oc_interno']
    insight = matriz.get('insight')
    if matriz['filas'] and matriz['columnas']:
        parrafo_matriz = (
            'Cruce entre el tipo de OC y el tipo interno de clasificación, solo para las OC fuera del PAC — '
            'ayuda a identificar en qué combinación se concentran las oportunidades de corrección.'
        )
        if insight:
            parrafo_matriz += (
                f" El cruce con mayor concentración es {insight['tipo_oc']} × {insight['tipo_interno']}, con "
                f"{_n(insight['cantidad'])} OC ({insight['pct_del_total']}% del total fuera del PAC analizado en "
                f"esta matriz) — es el punto de mayor impacto para priorizar revisión y posible enlace manual."
            )
        secciones.append(Seccion(
            titulo=f'Matriz Cruzada: Tipo OC × Tipo Interno ({anio})',
            nivel=2,
            parrafo=parrafo_matriz,
            tabla=Tabla(
                encabezados=['Tipo OC'] + matriz['columnas'] + ['Total'],
                filas=[
                    [fila] + [_n(matriz['datos'].get(fila, {}).get(col, 0)) for col in matriz['columnas']]
                    + [_n(sum(matriz['datos'].get(fila, {}).values()))]
                    for fila in matriz['filas']
                ],
            ),
        ))

    if corregidas.get('oc_unicas_corregidas', 0) > 0:
        secciones.append(Seccion(
            titulo='Corregidas — Revisiones Manuales de Enlace PAC',
            nivel=2,
            parrafo=(
                f"{_n(corregidas['oc_unicas_corregidas'])} OC fueron revisadas y enlazadas manualmente al PAC "
                f"(de {_n(corregidas['total_revisiones'])} revisiones registradas en el sistema)."
            ),
            tabla=Tabla(
                encabezados=['Concepto', 'Cantidad'],
                filas=[
                    ['Revisiones registradas', _n(corregidas['total_revisiones'])],
                    ['OC únicas corregidas', _n(corregidas['oc_unicas_corregidas'])],
                    ['Ya sincronizadas por el ETL', _n(corregidas['sincronizadas'])],
                    ['Esperando próxima sincronización', _n(corregidas['esperando_sync'])],
                ],
            ),
        ))

    secciones.append(Seccion(
        titulo='Conclusiones y Recomendaciones',
        parrafo=_parrafo_conclusiones(datos),
    ))

    return secciones


# =============================================================================
# Word (python-docx)
# =============================================================================

def _portada_ind1_docx(doc, anio, hoy):
    logo = _logo_bytes()
    if logo:
        p_logo = doc.add_paragraph()
        p_logo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_logo.add_run().add_picture(logo, width=Inches(1.5))

    doc.add_paragraph()
    titulo = doc.add_paragraph()
    titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_t = titulo.add_run(TITULO_INFORME_IND1)
    run_t.font.size = Pt(22)
    run_t.font.bold = True
    run_t.font.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_s = sub.add_run(NOMBRE_INSTITUCION.upper())
    run_s.font.size = Pt(15)
    run_s.font.bold = True
    run_s.font.color.rgb = RGBColor(0x38, 0xb2, 0xbd)

    edificio = _edificio_bytes()
    if edificio:
        doc.add_paragraph()
        p_img = doc.add_paragraph()
        p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_img.add_run().add_picture(edificio, width=Inches(5.6))

    doc.add_paragraph()
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_m = meta.add_run(f'Año PAC (Plan Anual de Compras): {anio}\nGenerado el {hoy.strftime("%d-%m-%Y")}')
    run_m.font.size = Pt(11.5)
    run_m.font.color.rgb = RGBColor(0x47, 0x55, 0x69)
    doc.add_page_break()


def _indice_ind1_docx(doc):
    h = doc.add_heading('Índice', level=1)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)
    p_ayuda = doc.add_paragraph()
    run_ayuda = p_ayuda.add_run('(el índice se actualiza automáticamente al abrir el documento; si no se ve, haga clic derecho sobre él y seleccione "Actualizar campos")')
    run_ayuda.font.size = Pt(9)
    run_ayuda.italic = True
    run_ayuda.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
    p_toc = doc.add_paragraph()
    _agregar_campo_docx(p_toc, 'TOC \\o "1-2" \\h \\z \\u')
    doc.add_page_break()


def _pie_pagina_ind1_docx(doc):
    section = doc.sections[0]
    footer = section.footer
    footer.is_linked_to_previous = True
    p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    color_pie = RGBColor(0x94, 0xa3, 0xb8)
    run1 = p.add_run(f'{NOMBRE_INSTITUCION} · Indicador 1 Res.188/2026 · Página ')
    run1.font.size = Pt(8)
    run1.font.color.rgb = color_pie
    _agregar_campo_docx(p, 'PAGE')
    run2 = p.add_run(' de ')
    run2.font.size = Pt(8)
    run2.font.color.rgb = color_pie
    _agregar_campo_docx(p, 'NUMPAGES')


def _tabla_docx(doc, tabla: Tabla):
    n_cols = len(tabla.encabezados)
    t = doc.add_table(rows=1, cols=n_cols)
    t.style = 'Light Grid Accent 1'
    for i, h in enumerate(tabla.encabezados):
        t.rows[0].cells[i].text = str(h)
    for fila in tabla.filas:
        celdas = t.add_row().cells
        for i, val in enumerate(fila):
            celdas[i].text = str(val)
    if tabla.anchos_word:
        for i, ancho in enumerate(tabla.anchos_word):
            for row in t.rows:
                row.cells[i].width = Inches(ancho)


def _render_seccion_docx(doc, s: Seccion):
    if s.nivel == 1:
        _titulo_capitulo_docx(doc, s.titulo)
    else:
        _titulo_seccion_docx(doc, s.titulo)
    if s.parrafo:
        doc.add_paragraph(s.parrafo)
    for g in s.graficos:
        if not g.imagen:
            continue
        doc.add_picture(g.imagen, width=Inches(g.ancho_word))
        cap = doc.add_paragraph()
        run = cap.add_run(g.caption)
        run.italic = True
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
    if s.tabla:
        _tabla_docx(doc, s.tabla)
        doc.add_paragraph()
    if s.nota:
        p = doc.add_paragraph()
        run = p.add_run('📋 ' + s.nota)
        run.italic = True
        run.font.size = Pt(9.5)


def generar_informe_word_ind1(anio):
    """Genera el informe Word del Indicador 1 (% Compras dentro del PAC). Retorna BytesIO."""
    datos = _construir_datos(anio)
    secciones = _construir_secciones(datos)

    doc = Document()
    _configurar_margenes_docx(doc)
    _portada_ind1_docx(doc, anio, datos['hoy'])
    _indice_ind1_docx(doc)
    for s in secciones:
        _render_seccion_docx(doc, s)
    _pie_pagina_ind1_docx(doc)
    _forzar_actualizacion_campos_docx(doc)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# =============================================================================
# PDF (reportlab)
# =============================================================================

class _InformeInd1DocTemplate(BaseDocTemplate):
    """Mismo patrón que `_InformeDocTemplate` de services_reportes.py (pie de
    página + registro de entradas de índice vía `notify('TOCEntry', ...)`),
    con texto de pie propio — por eso no se reutiliza la clase original."""

    def __init__(self, *args, **kwargs):
        BaseDocTemplate.__init__(self, *args, **kwargs)
        frame = Frame(self.leftMargin, self.bottomMargin, self.width, self.height, id='normal')
        self.addPageTemplates([PageTemplate(id='con_pie', frames=[frame], onPage=self._dibujar_pie)])

    def _dibujar_pie(self, canvas, doc):
        canvas.saveState()
        canvas.setFont('Helvetica', 7.5)
        canvas.setFillColor(colors.HexColor('#94a3b8'))
        canvas.drawCentredString(doc.pagesize[0] / 2, 0.45 * inch, f'{NOMBRE_INSTITUCION} · Indicador 1 Res.188/2026 · Página {doc.page}')
        canvas.restoreState()

    def afterFlowable(self, flowable):
        if isinstance(flowable, Paragraph):
            estilo = flowable.style.name
            texto = flowable.getPlainText()
            if estilo == 'TituloCapitulo':
                self.notify('TOCEntry', (0, texto, self.page))
            elif estilo == 'TituloSeccion':
                self.notify('TOCEntry', (1, texto, self.page))


def _portada_ind1_pdf(story, anio, hoy):
    logo = _pdf_imagen(_logo_bytes(), 1.3, 1.3)
    if logo:
        story.append(logo)
    story.append(Spacer(1, 0.25 * inch))
    story.append(Paragraph(TITULO_INFORME_IND1, _PDF_ESTILOS['TituloPortada']))
    story.append(Paragraph(NOMBRE_INSTITUCION.upper(), _PDF_ESTILOS['SubportadaFuerte']))
    edificio = _pdf_imagen(_edificio_bytes(), 5.3, 3.48)
    if edificio:
        story.append(Spacer(1, 0.2 * inch))
        story.append(edificio)
    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph(f'Año PAC (Plan Anual de Compras): {anio}<br/>Generado el {hoy.strftime("%d-%m-%Y")}', _PDF_ESTILOS['Subportada']))
    story.append(PageBreak())


def _indice_ind1_pdf(story):
    toc = TableOfContents()
    toc.levelStyles = [
        ParagraphStyle(name='TOCCapitulo', fontSize=11, leading=16, textColor=colors.HexColor(COLOR_INSTITUCIONAL), spaceBefore=4),
        ParagraphStyle(name='TOCSeccion', fontSize=9.5, leading=13, leftIndent=16, textColor=colors.HexColor('#475569')),
    ]
    story.append(Paragraph('Índice', _PDF_ESTILOS['TituloCapitulo']))
    story.append(toc)
    story.append(PageBreak())


def _tabla_pdf(tabla: Tabla):
    filas = [_fila_encabezado(tabla.encabezados)] + [[_celda(v) for v in fila] for fila in tabla.filas]
    n_cols = len(tabla.encabezados)
    ancho_disponible = 6.3
    anchos = tabla.anchos_pdf or [ancho_disponible / n_cols] * n_cols
    t = Table(filas, colWidths=[a * inch for a in anchos], repeatRows=1)
    t.setStyle(_PDF_TABLA_ESTILO)
    return t


def _render_seccion_pdf(story, s: Seccion):
    story.append(Paragraph(s.titulo, _PDF_ESTILOS['TituloCapitulo' if s.nivel == 1 else 'TituloSeccion']))
    if s.parrafo:
        story.append(Paragraph(s.parrafo, _PDF_ESTILOS['Cuerpo']))
    for g in s.graficos:
        img = _pdf_imagen(g.imagen, *g.ancho_pdf)
        if img:
            story.append(img)
            story.append(Paragraph(g.caption, _PDF_ESTILOS['Leyenda']))
    if s.tabla:
        story.append(_tabla_pdf(s.tabla))
        story.append(Spacer(1, 0.15 * inch))
    if s.nota:
        story.append(Paragraph('📋 ' + s.nota, _PDF_ESTILOS['Leyenda']))
    story.append(Spacer(1, 0.1 * inch))


def generar_reporte_pdf_ind1(anio):
    """Genera el informe PDF del Indicador 1. Retorna BytesIO. Misma estructura que `generar_informe_word_ind1`."""
    datos = _construir_datos(anio)
    secciones = _construir_secciones(datos)

    buf = io.BytesIO()
    doc = _InformeInd1DocTemplate(buf, pagesize=letter, topMargin=0.9 * inch, bottomMargin=0.85 * inch, leftMargin=0.95 * inch, rightMargin=0.75 * inch)
    story = []
    _portada_ind1_pdf(story, anio, datos['hoy'])
    _indice_ind1_pdf(story)
    for s in secciones:
        _render_seccion_pdf(story, s)

    doc.multiBuild(story)
    buf.seek(0)
    return buf


# =============================================================================
# PPT (python-pptx) — versión ejecutiva, un slide por sección
# =============================================================================

def _ppt_tabla_generica(slide, tabla: Tabla, top, limite=8):
    filas_datos = tabla.filas[:limite]
    n_filas = len(filas_datos) + 1
    n_cols = len(tabla.encabezados)
    tabla_shape = slide.shapes.add_table(n_filas, n_cols, PptxInches(0.5), top, PptxInches(9), PptxInches(0.4 * n_filas))
    t = tabla_shape.table
    for i, h in enumerate(tabla.encabezados):
        t.cell(0, i).text = str(h)
    for r, fila in enumerate(filas_datos, start=1):
        for c, val in enumerate(fila):
            t.cell(r, c).text = str(val)
    for row in t.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = PptxPt(11)


def _render_seccion_ppt(prs, s: Seccion):
    slide = _ppt_slide_en_blanco(prs)
    _ppt_titulo(slide, s.titulo, tamano=22 if s.nivel == 1 else 19)
    top = PptxInches(1.1)
    if s.parrafo:
        _ppt_parrafo(slide, s.parrafo, top, height=PptxInches(0.9), tamano=13)
        top = PptxInches(1.85)
    if s.graficos and s.graficos[0].imagen:
        g = s.graficos[0]
        _ppt_imagen(slide, g.imagen, PptxInches((10 - g.ancho_ppt) / 2), top, PptxInches(g.ancho_ppt))
    elif s.tabla:
        _ppt_tabla_generica(slide, s.tabla, top)
    _ppt_pie_pagina(slide)
    return slide


def generar_presentacion_ppt_ind1(anio):
    """Genera la presentación PPT del Indicador 1 — un slide por sección (versión
    ejecutiva). Retorna BytesIO."""
    datos = _construir_datos(anio)
    secciones = _construir_secciones(datos)

    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)

    slide = _ppt_slide_en_blanco(prs)
    _ppt_titulo(slide, TITULO_INFORME_IND1, top=PptxInches(2.6), tamano=28)
    _ppt_parrafo(slide, f'{NOMBRE_INSTITUCION}\nAño PAC: {anio} · Generado el {datos["hoy"].strftime("%d-%m-%Y")}', PptxInches(3.6), tamano=14)
    _ppt_pie_pagina(slide)

    for s in secciones:
        _render_seccion_ppt(prs, s)
        if s.tabla and s.graficos and s.graficos[0].imagen:
            # Sección con gráfico Y tabla: la tabla va en un segundo slide (versión ejecutiva no las mezcla).
            slide2 = _ppt_slide_en_blanco(prs)
            _ppt_titulo(slide2, s.titulo + ' — Detalle', tamano=19)
            _ppt_tabla_generica(slide2, s.tabla, PptxInches(1.1))
            _ppt_pie_pagina(slide2)

    _ppt_numerar_diapositivas(prs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
